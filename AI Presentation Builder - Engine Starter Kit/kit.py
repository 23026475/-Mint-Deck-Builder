from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

NAVY=RGBColor(0x24,0x3C,0x7A); NAVY2=RGBColor(0x18,0x4B,0x9A); DARK=RGBColor(0x1B,0x2B,0x4E)
CYAN=RGBColor(0x1B,0xA9,0xC4); GREEN=RGBColor(0x43,0x9E,0x46); LGREEN=RGBColor(0x89,0xC1,0x46)
ORANGE=RGBColor(0xF1,0x59,0x27); GREY=RGBColor(0x47,0x50,0x54); MGREY=RGBColor(0x6C,0x77,0x7D)
WHITE=RGBColor(0xFF,0xFF,0xFF); BORDER=RGBColor(0xE3,0xE7,0xEC); PEACH=RGBColor(0xFD,0xEE,0xE4); DORANGE=RGBColor(0xB5,0x45,0x0F)
CARD_COLORS=[CYAN,NAVY,ORANGE,GREEN]
SW,SH=Emu(12192000),Emu(6858000)
IN=914400

def _noline(sh): sh.line.fill.background()
def _shadow_off(sh):
    try: sh.shadow.inherit=False
    except: pass
def rect(s,x,y,w,h,color,shape=MSO_SHAPE.RECTANGLE,line=None):
    sh=s.shapes.add_shape(shape,Emu(int(x)),Emu(int(y)),Emu(int(w)),Emu(int(h)))
    sh.fill.solid(); sh.fill.fore_color.rgb=color; _shadow_off(sh)
    if line: sh.line.color.rgb=line; sh.line.width=Pt(0.75)
    else: _noline(sh)
    return sh
def txt(s,x,y,w,h,runs,size=12,color=GREY,bold=False,italic=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,spacing=None,wrap=True):
    tb=s.shapes.add_textbox(Emu(int(x)),Emu(int(y)),Emu(int(w)),Emu(int(h)))
    tf=tb.text_frame; tf.word_wrap=wrap; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    if isinstance(runs,str): runs=[(runs,{})]
    first=True
    for text,ov in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        if spacing: p.line_spacing=spacing
        if ov.get('space_before'): p.space_before=Pt(ov['space_before'])
        p.alignment=ov.get('align',align)
        r=p.add_run(); r.text=text
        f=r.font; f.size=Pt(ov.get('size',size)); f.bold=ov.get('bold',bold)
        f.italic=ov.get('italic',italic); f.color.rgb=ov.get('color',color)
        f.name='Lato'
    return tb
def base(prs,layout):
    s=prs.slides.add_slide(layout)
    rect(s,0,0,SW,SH,WHITE)
    return s
def footer(s,left,page):
    txt(s,0.45*IN,6.62*IN,9.5*IN,0.3*IN,left,size=8.5,color=MGREY)
    txt(s,12.5*IN,6.62*IN,0.5*IN,0.3*IN,f'{page:02d}',size=8.5,color=MGREY,align=PP_ALIGN.RIGHT)
def header(s,kicker,title,intro=None,kcolor=CYAN):
    txt(s,0.45*IN,0.32*IN,10*IN,0.3*IN,kicker.upper(),size=10.5,color=kcolor,bold=True)
    txt(s,0.43*IN,0.62*IN,12.4*IN,0.6*IN,title,size=23,color=RGBColor(0x2A,0x33,0x38),bold=True)
    if intro:
        txt(s,0.45*IN,1.28*IN,12.4*IN,0.6*IN,intro,size=11,color=GREY,spacing=1.15)
def card(s,x,y,w,h,color,title,body=None,bullets=None,tsize=12.5,bsize=9.5):
    c=rect(s,x,y,w,h,WHITE,MSO_SHAPE.ROUNDED_RECTANGLE,line=BORDER)
    c.adjustments[0]=0.045
    rect(s,x+0.06*IN,y,w-0.12*IN,0.075*IN,color)
    txt(s,x+0.18*IN,y+0.22*IN,w-0.36*IN,0.35*IN,title,size=tsize,color=color,bold=True)
    ty=y+0.62*IN
    if body:
        txt(s,x+0.18*IN,ty,w-0.36*IN,h-(ty-y)-0.15*IN,body,size=bsize,color=GREY,spacing=1.12)
    if bullets:
        runs=[('•  '+b,{'space_before':4}) for b in bullets]
        txt(s,x+0.18*IN,ty,w-0.36*IN,h-(ty-y)-0.15*IN,runs,size=bsize,color=GREY,spacing=1.05)
def card_row(s,items,y=2.0*IN,h=2.6*IN,x0=0.45*IN,gap=0.22*IN,bullets=False):
    n=len(items)
    w=(12.6*IN-2*x0-(n-1)*gap)/n
    for i,it in enumerate(items):
        col=CARD_COLORS[i%4]
        if bullets: card(s,x0+i*(w+gap),y,w,h,col,it[0],bullets=it[1])
        else: card(s,x0+i*(w+gap),y,w,h,col,it[0],body=it[1])
def callout(s,y,lead,body,tint=PEACH,lead_color=ORANGE,h=0.85*IN):
    b=rect(s,0.45*IN,y,12.45*IN,h,tint,MSO_SHAPE.ROUNDED_RECTANGLE); b.adjustments[0]=0.12
    txt(s,0.75*IN,y+0.14*IN,11.9*IN,h-0.28*IN,[(lead+'  ',{'bold':True,'color':lead_color}),],size=10.5,color=GREY)
    # lead + body on same frame: rebuild with two runs in one paragraph
def callout2(s,y,lead,body,tint=PEACH,lead_color=DORANGE,h=0.9*IN):
    b=rect(s,0.45*IN,y,12.45*IN,h,tint,MSO_SHAPE.ROUNDED_RECTANGLE); b.adjustments[0]=0.12
    tb=s.shapes.add_textbox(Emu(int(0.75*IN)),Emu(int(y+0.13*IN)),Emu(int(11.9*IN)),Emu(int(h-0.26*IN)))
    tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    p=tf.paragraphs[0]; p.line_spacing=1.15
    r1=p.add_run(); r1.text=lead+'  '; r1.font.bold=True; r1.font.color.rgb=lead_color; r1.font.size=Pt(10.5); r1.font.name='Lato'
    r2=p.add_run(); r2.text=body; r2.font.color.rgb=GREY; r2.font.size=Pt(10.5); r2.font.name='Lato'
def flow_row(s,steps,y=2.1*IN,h=1.8*IN,x0=0.45*IN):
    n=len(steps); gap=0.5*IN
    w=(12.6*IN-2*x0-(n-1)*gap)/n
    for i,(t,b) in enumerate(steps):
        x=x0+i*(w+gap)
        card(s,x,y,w,h,CARD_COLORS[i%4],t,body=b)
        if i<n-1:
            txt(s,x+w+0.08*IN,y+h/2-0.2*IN,gap-0.14*IN,0.4*IN,'→',size=18,color=CYAN,bold=True,align=PP_ALIGN.CENTER)
def statcard(s,x,y,w,h,kicker,big,sub,foot1,foot2=None):
    rect(s,x,y+0.06*IN,w,h-0.06*IN,DARK,MSO_SHAPE.ROUNDED_RECTANGLE).adjustments[0]=0.05
    rect(s,x+0.06*IN,y,w-0.12*IN,0.09*IN,LGREEN)
    txt(s,x+0.3*IN,y+0.35*IN,w-0.6*IN,0.3*IN,kicker.upper(),size=10.5,color=CYAN,bold=True)
    txt(s,x+0.3*IN,y+0.75*IN,w-0.6*IN,0.7*IN,big,size=27,color=WHITE,bold=True)
    txt(s,x+0.3*IN,y+1.42*IN,w-0.6*IN,0.35*IN,sub,size=13,color=LGREEN,bold=True)
    txt(s,x+0.3*IN,y+h-0.85*IN,w-0.6*IN,0.35*IN,foot1,size=10.5,color=WHITE,bold=True)
    if foot2: txt(s,x+0.3*IN,y+h-0.52*IN,w-0.6*IN,0.4*IN,foot2,size=9,color=RGBColor(0xB9,0xC4,0xD6),italic=True)
def bullets_col(s,x,y,w,items,size=10,color=GREY,dot=ORANGE,gap=6):
    runs=[]
    for it in items: runs.append(('•  '+it,{'space_before':gap}))
    tb=txt(s,x,y,w,4.5*IN,runs,size=size,color=color,spacing=1.1)
    return tb
def cover(prs,layout,kicker,title,sub,abstract,foot):
    s=prs.slides.add_slide(layout)
    rect(s,0,0,SW,SH,DARK)
    rect(s,0,0,SW,0.09*IN,CYAN); rect(s,0,SH-0.09*IN,SW,0.09*IN,LGREEN)
    txt(s,0.6*IN,1.9*IN,7*IN,0.35*IN,kicker.upper(),size=12,color=CYAN,bold=True)
    txt(s,0.58*IN,2.3*IN,7.4*IN,1.0*IN,title,size=40,color=WHITE,bold=True)
    txt(s,0.6*IN,3.35*IN,7.4*IN,0.5*IN,sub,size=19,color=RGBColor(0xBD,0xD3,0xF0),bold=True)
    txt(s,0.6*IN,4.35*IN,6.9*IN,0.9*IN,abstract,size=11.5,color=RGBColor(0xC9,0xD2,0xE4),italic=True,spacing=1.25)
    txt(s,0.6*IN,6.1*IN,8*IN,0.3*IN,foot,size=10,color=CYAN,bold=True)
    return s
def closing(prs,layout,kicker,title,steps,foot):
    s=prs.slides.add_slide(layout)
    rect(s,0,0,SW,SH,DARK)
    rect(s,0,0,SW,0.09*IN,CYAN); rect(s,0,SH-0.09*IN,SW,0.09*IN,LGREEN)
    txt(s,0.6*IN,1.5*IN,8*IN,0.35*IN,kicker.upper(),size=12,color=CYAN,bold=True)
    txt(s,0.58*IN,1.9*IN,11.9*IN,0.8*IN,title,size=30,color=WHITE,bold=True)
    runs=[]
    for i,st in enumerate(steps,1):
        runs.append((f'{i}.   {st}',{'space_before':14}))
    txt(s,0.62*IN,3.0*IN,11.5*IN,2.6*IN,runs,size=13,color=RGBColor(0xE6,0xEB,0xF4),spacing=1.15)
    txt(s,0.6*IN,6.05*IN,9*IN,0.3*IN,foot,size=10,color=CYAN,bold=True)
    return s
def get_blank(prs):
    for l in prs.slide_masters[0].slide_layouts:
        if l.name=='Blank': return l
    return prs.slide_masters[0].slide_layouts[0]
