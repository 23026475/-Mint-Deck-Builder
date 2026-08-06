"""Integration tests for the first milestone DeckBuilder pipeline."""

from __future__ import annotations

import json
import tempfile
import unittest
from pathlib import Path
from types import SimpleNamespace

from pptx import Presentation

from presentation_engine.builders.deck_builder import DeckBuilder


class FakeTemplateRetriever:
    def __init__(self, potx_path: Path) -> None:
        self.potx_path = potx_path

    def retrieve_latest_template(self) -> Path:
        return self.potx_path


class FakeTemplateConverter:
    def __init__(self, converted_pptx_path: Path) -> None:
        self.converted_pptx_path = converted_pptx_path
        self.calls: list[tuple[Path, Path, bool]] = []

    def convert(self, potx_path: str | Path, pptx_path: str | Path, overwrite: bool = True) -> Path:
        self.calls.append((Path(potx_path), Path(pptx_path), overwrite))
        return self.converted_pptx_path


class NoOpPlaceholderBuilder:
    """
    Keeps this integration test focused on DeckBuilder orchestration.

    PlaceholderBuilder has its own dedicated unit tests, so this test should not
    depend on real FY27 placeholder indexes.
    """

    def populate(self, slide, slide_definition) -> None:
        return None


class DeckBuilderIntegrationTests(unittest.TestCase):
    def test_generates_presentation_with_exactly_milestone_slides(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            base_dir = Path(temp_dir)
            config = self._create_config(base_dir)
            config.ensure_runtime_directories()

            self._write_archetype_baseline(base_dir)
            self._write_sample_contract(config.input.contract_path)

            potx_path = config.template.downloaded_potx_path
            potx_path.write_bytes(b"local development potx placeholder")

            self._create_converted_template_with_sample_slides(
                config.template.converted_pptx_path
            )

            builder = DeckBuilder(
                engine_config=config,
                template_retriever=FakeTemplateRetriever(potx_path),
                template_converter=FakeTemplateConverter(config.template.converted_pptx_path),
                placeholder_builder=NoOpPlaceholderBuilder(),
            )

            result = builder.build_from_contract_file()

            self.assertTrue(result.output_pptx_path.exists())

            generated = Presentation(str(result.output_pptx_path))

            self.assertEqual(len(generated.slides), 3)
            self.assertEqual(result.slide_count, 3)

    def _create_config(self, base_dir: Path) -> SimpleNamespace:
        work_dir = base_dir / "data" / "work"
        template_dir = work_dir / "template"
        input_dir = base_dir / "data" / "input"
        output_dir = base_dir / "data" / "output"
        log_dir = base_dir / "logs"

        config = SimpleNamespace(
            base_dir=base_dir,
            work_dir=work_dir,
            template=SimpleNamespace(
                download_dir=template_dir,
                downloaded_potx_path=template_dir / "FY27 AI-Ready v3.0.potx",
                converted_pptx_path=template_dir / "FY27 AI-Ready v3.0.work.pptx",
            ),
            input=SimpleNamespace(
                input_dir=input_dir,
                contract_path=input_dir / "sample_deck_contract.json",
            ),
            output=SimpleNamespace(output_dir=output_dir),
            logging=SimpleNamespace(log_dir=log_dir),
        )

        def ensure_runtime_directories() -> None:
            for directory in [work_dir, template_dir, input_dir, output_dir, log_dir]:
                directory.mkdir(parents=True, exist_ok=True)

        config.ensure_runtime_directories = ensure_runtime_directories
        return config

    def _write_archetype_baseline(self, base_dir: Path) -> None:
        config_dir = base_dir / "config"
        config_dir.mkdir(parents=True, exist_ok=True)

        baseline = {
            "archetypes": [
                {"archetype": "cover", "layout": "Blank"},
                {"archetype": "cards3", "layout": "Blank"},
                {"archetype": "closing", "layout": "Blank"},
            ]
        }

        (config_dir / "archetype-baseline.json").write_text(
            json.dumps(baseline),
            encoding="utf8",
        )

    def _write_sample_contract(self, contract_path: Path) -> None:
        contract = {
            "deck": {
                "client": "Acme Bank",
                "title": "Milestone Test",
                "mode": "STAGE",
            },
            "slides": [
                {
                    "archetype": "cover",
                    "fields": {
                        "title": "Milestone Test",
                    },
                },
                {
                    "archetype": "cards3",
                    "action_title": "Three risks need attention",
                    "fields": {},
                },
                {
                    "archetype": "closing",
                    "action_title": "From today to kickoff",
                    "fields": {},
                },
            ],
        }

        contract_path.write_text(json.dumps(contract), encoding="utf8")

    def _create_converted_template_with_sample_slides(self, template_path: Path) -> None:
        presentation = Presentation()
        blank_layout = presentation.slide_layouts[6]

        presentation.slides.add_slide(blank_layout)
        presentation.slides.add_slide(blank_layout)
        presentation.slides.add_slide(blank_layout)

        template_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(template_path))


if __name__ == "__main__":
    unittest.main()