"""Identify the exact geometry assignment that writes cx=0.

Diagnostic only. This script does not modify production source files.

It runs the existing engine using the real IT Skills contract and instruments
TextBoxLayoutPolicy.apply() in memory. For each target placeholder it:

1. Captures the original inherited state.
2. Executes the current full sequence one assignment at a time:
   left, top, width, height.
3. Restores the original placeholder XML.
4. Executes a height-only control.
5. Restores the original placeholder XML again.
6. Calls the real policy so the normal build can finish.

Intermediate observations are captured before presentation save().

Run from the project root:

    $env:PYTHONPATH = ".\\src"
    python .\\diagnose_textbox_geometry_assignments.py

Outputs:

    data/output/textbox_geometry_assignment_diagnostic.json
    data/output/textbox_geometry_assignment_diagnostic.md
"""

from __future__ import annotations

import argparse
import copy
import json
import math
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Any, Mapping

from pptx.oxml.ns import qn
from pptx.util import Inches

from presentation_engine.builders.deck_builder import DeckBuilder
from presentation_engine.services.text_box_layout import TextBoxLayoutPolicy


EMU_PER_INCH = 914400
TARGETS = {
    ("statement", "support"): {"slide_number": 2, "idx": 2},
    ("image_right", "body"): {"slide_number": 6, "idx": 12},
    ("closing", "steps"): {"slide_number": 8, "idx": 12},
}


@dataclass
class MutationSnapshot:
    experiment: str
    operation: str
    slide_number: int
    archetype: str
    placeholder_name: str
    idx: int
    width_value: int | None
    width_type: str
    height_value: int | None
    left_value: int | None
    top_value: int | None
    width_inches: float | None
    height_inches: float | None
    left_inches: float | None
    top_inches: float | None
    xfrm_exists: bool
    off_x: int | None
    off_y: int | None
    ext_cx: int | None
    ext_cy: int | None
    xfrm_xml: str | None
    sppr_xml: str | None
    layout_width_value: int | None
    layout_width_type: str
    layout_width_inches: float | None
    layout_ext_cx: int | None
    layout_xfrm_xml: str | None


def as_int(value: Any) -> int | None:
    if value is None:
        return None
    try:
        return int(value)
    except (TypeError, ValueError):
        return None


def as_inches(value: Any) -> float | None:
    integer = as_int(value)
    return None if integer is None else round(integer / EMU_PER_INCH, 4)


def placeholder_idx(shape: Any) -> int | None:
    try:
        return int(shape.placeholder_format.idx)
    except (AttributeError, TypeError, ValueError):
        return None


def matching_layout_placeholder(shape: Any) -> Any | None:
    slide = getattr(getattr(shape, "part", None), "slide", None)
    if slide is None:
        return None
    idx = placeholder_idx(shape)
    for candidate in slide.slide_layout.placeholders:
        if placeholder_idx(candidate) == idx:
            return candidate
    return None


def sppr(shape: Any) -> Any:
    return shape._element.spPr


def xfrm(shape: Any) -> Any | None:
    return getattr(sppr(shape), "xfrm", None)


def child_value(parent: Any, tag: str, attribute: str) -> int | None:
    if parent is None:
        return None
    child = parent.find(qn(tag))
    if child is None:
        return None
    return as_int(child.get(attribute))


def xml(element: Any) -> str | None:
    if element is None:
        return None
    try:
        return element.xml
    except AttributeError:
        return str(element)


def capture(
    *,
    experiment: str,
    operation: str,
    shape: Any,
    archetype: str,
    placeholder_name: str,
    slide_number: int,
    idx: int,
) -> MutationSnapshot:
    layout = matching_layout_placeholder(shape)
    shape_xfrm = xfrm(shape)
    layout_xfrm = xfrm(layout) if layout is not None else None

    width = getattr(shape, "width", None)
    layout_width = getattr(layout, "width", None) if layout is not None else None

    return MutationSnapshot(
        experiment=experiment,
        operation=operation,
        slide_number=slide_number,
        archetype=archetype,
        placeholder_name=placeholder_name,
        idx=idx,
        width_value=as_int(width),
        width_type=type(width).__name__,
        height_value=as_int(getattr(shape, "height", None)),
        left_value=as_int(getattr(shape, "left", None)),
        top_value=as_int(getattr(shape, "top", None)),
        width_inches=as_inches(width),
        height_inches=as_inches(getattr(shape, "height", None)),
        left_inches=as_inches(getattr(shape, "left", None)),
        top_inches=as_inches(getattr(shape, "top", None)),
        xfrm_exists=shape_xfrm is not None,
        off_x=child_value(shape_xfrm, "a:off", "x"),
        off_y=child_value(shape_xfrm, "a:off", "y"),
        ext_cx=child_value(shape_xfrm, "a:ext", "cx"),
        ext_cy=child_value(shape_xfrm, "a:ext", "cy"),
        xfrm_xml=xml(shape_xfrm),
        sppr_xml=xml(sppr(shape)),
        layout_width_value=as_int(layout_width),
        layout_width_type=type(layout_width).__name__,
        layout_width_inches=as_inches(layout_width),
        layout_ext_cx=child_value(layout_xfrm, "a:ext", "cx"),
        layout_xfrm_xml=xml(layout_xfrm),
    )


def restore_sppr(shape: Any, original_sppr: Any) -> None:
    element = shape._element
    current = element.spPr
    restored = copy.deepcopy(original_sppr)
    element.replace(current, restored)


def configured_height(
    policy: TextBoxLayoutPolicy,
    *,
    shape: Any,
    archetype: str,
    placeholder_name: str,
    text: str,
    font_size_pt: float | None,
) -> float:
    """Reproduce only the current policy's height calculation, without writes."""
    rule = policy.resolve(archetype, placeholder_name)
    if rule is None:
        raise RuntimeError(f"No rule for {archetype}.{placeholder_name}")

    width = getattr(shape, "width", None)
    width_inches = float(width) / EMU_PER_INCH
    resolved_font_size = font_size_pt or font_size_from_shape(shape) or 14.0
    defaults: Mapping[str, Any] = policy._config["defaults"]

    line_height_multiplier = float(defaults.get("line_height_multiplier", 1.2))
    paragraph_gap_lines = float(defaults.get("paragraph_gap_lines", 0.35))
    characters_per_point_inch = float(
        defaults.get("characters_per_point_inch", 27.8)
    )
    margin_left = float(defaults.get("margin_left", 0.08))
    margin_right = float(defaults.get("margin_right", 0.08))
    margin_top = float(defaults.get("margin_top", 0.05))
    margin_bottom = float(defaults.get("margin_bottom", 0.05))

    usable_width = max(0.1, width_inches - margin_left - margin_right)
    chars_per_inch = max(
        8.0,
        characters_per_point_inch * 14.0 / resolved_font_size,
    )
    chars_per_line = max(1, int(usable_width * chars_per_inch))

    paragraphs = (text or "").splitlines() or [""]
    estimated_lines = 0
    for paragraph in paragraphs:
        normalized = " ".join(paragraph.split())
        estimated_lines += max(1, math.ceil(len(normalized) / chars_per_line))
    if len(paragraphs) > 1:
        estimated_lines += math.ceil((len(paragraphs) - 1) * paragraph_gap_lines)

    line_height_inches = (resolved_font_size / 72.0) * line_height_multiplier
    height = estimated_lines * line_height_inches + margin_top + margin_bottom
    if rule.min_height is not None:
        height = max(height, rule.min_height)
    if rule.max_height is not None:
        height = min(height, rule.max_height)
    return height


def font_size_from_shape(shape: Any) -> float | None:
    text_frame = getattr(shape, "text_frame", None)
    if text_frame is None:
        return None
    for paragraph in list(getattr(text_frame, "paragraphs", []) or []):
        paragraph_font = getattr(paragraph, "font", None)
        size = getattr(paragraph_font, "size", None) if paragraph_font else None
        if size is not None:
            return float(size.pt)
        for run in list(getattr(paragraph, "runs", []) or []):
            run_font = getattr(run, "font", None)
            size = getattr(run_font, "size", None) if run_font else None
            if size is not None:
                return float(size.pt)
    return None


def first_zero_transition(items: list[MutationSnapshot]) -> dict[str, str]:
    conclusions: dict[str, str] = {}
    grouped: dict[tuple[int, str], list[MutationSnapshot]] = {}
    for item in items:
        grouped.setdefault((item.slide_number, item.experiment), []).append(item)

    for (slide, experiment), states in grouped.items():
        conclusion = "No cx=0 transition observed."
        previous: MutationSnapshot | None = None
        for current in states:
            if previous is not None:
                previous_cx = previous.ext_cx
                current_cx = current.ext_cx
                previous_width = previous.width_value
                current_width = current.width_value
                if (
                    (previous_cx is None or previous_cx > 0 or (previous_width or 0) > 0)
                    and (current_cx == 0 or current_width == 0)
                ):
                    conclusion = (
                        f"First zero-width transition occurs at operation "
                        f"'{current.operation}', following '{previous.operation}'."
                    )
                    break
            previous = current
        conclusions[f"slide_{slide}_{experiment}"] = conclusion
    return conclusions


def row(item: MutationSnapshot) -> str:
    def fmt(value: float | None) -> str:
        return "None" if value is None else f"{value:.2f}"
    return (
        f"| {item.slide_number} | {item.idx} | {item.experiment} | "
        f"{item.operation} | {fmt(item.width_inches)} | "
        f"{fmt(item.height_inches)} | {item.ext_cx} | {item.ext_cy} | "
        f"{item.width_type} | {item.xfrm_exists} |"
    )


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--contract",
        default="data/input/it_skills_engine_validation_contract.json",
    )
    parser.add_argument(
        "--json-output",
        default="data/output/textbox_geometry_assignment_diagnostic.json",
    )
    parser.add_argument(
        "--markdown-output",
        default="data/output/textbox_geometry_assignment_diagnostic.md",
    )
    args = parser.parse_args()

    snapshots: list[MutationSnapshot] = []
    original_apply = TextBoxLayoutPolicy.apply

    def diagnostic_apply(
        self: TextBoxLayoutPolicy,
        shape: Any,
        *,
        archetype: str,
        placeholder_name: str,
        text: str,
        font_size_pt: float | None = None,
    ) -> Any:
        target = TARGETS.get((archetype, placeholder_name))
        if target is None:
            return original_apply(
                self,
                shape,
                archetype=archetype,
                placeholder_name=placeholder_name,
                text=text,
                font_size_pt=font_size_pt,
            )

        original_sppr = copy.deepcopy(sppr(shape))
        left = getattr(shape, "left", None)
        top = getattr(shape, "top", None)
        width = getattr(shape, "width", None)
        applied_height = configured_height(
            self,
            shape=shape,
            archetype=archetype,
            placeholder_name=placeholder_name,
            text=text,
            font_size_pt=font_size_pt,
        )

        def snap(experiment: str, operation: str) -> None:
            snapshots.append(
                capture(
                    experiment=experiment,
                    operation=operation,
                    shape=shape,
                    archetype=archetype,
                    placeholder_name=placeholder_name,
                    slide_number=target["slide_number"],
                    idx=target["idx"],
                )
            )

        # Full current sequence, one assignment at a time.
        snap("full_sequence", "state_0_before_any_assignment")
        shape.left = left
        snap("full_sequence", "state_1_after_left")
        shape.top = top
        snap("full_sequence", "state_2_after_top")
        shape.width = width
        snap("full_sequence", "state_3_after_width")
        shape.height = Inches(applied_height)
        snap("full_sequence", "state_4_after_height")

        # Restore exact original spPr before the independent control.
        restore_sppr(shape, original_sppr)
        snap("height_only", "state_0_restored_before_height")
        shape.height = Inches(applied_height)
        snap("height_only", "state_1_after_height_only")

        # Restore exact original spPr before allowing the real policy to run.
        restore_sppr(shape, original_sppr)

        return original_apply(
            self,
            shape,
            archetype=archetype,
            placeholder_name=placeholder_name,
            text=text,
            font_size_pt=font_size_pt,
        )

    TextBoxLayoutPolicy.apply = diagnostic_apply
    try:
        DeckBuilder().build_from_contract_file(Path(args.contract))
    finally:
        TextBoxLayoutPolicy.apply = original_apply

    conclusions = first_zero_transition(snapshots)
    payload = {
        "contract": args.contract,
        "snapshots": [asdict(item) for item in snapshots],
        "conclusions": conclusions,
    }

    json_path = Path(args.json_output)
    md_path = Path(args.markdown_output)
    json_path.parent.mkdir(parents=True, exist_ok=True)
    md_path.parent.mkdir(parents=True, exist_ok=True)
    json_path.write_text(json.dumps(payload, indent=2), encoding="utf8")

    lines = [
        "# TextBox Geometry Assignment Diagnostic",
        "",
        "| Slide | idx | Experiment | Operation | Width (in) | Height (in) | cx | cy | width type | xfrm exists |",
        "|---:|---:|---|---|---:|---:|---:|---:|---|---|",
    ]
    lines.extend(row(item) for item in snapshots)
    lines.extend(["", "## Conclusions", ""])
    for key, value in conclusions.items():
        lines.append(f"- **{key}**: {value}")
    lines.extend(
        [
            "",
            "## Complete evidence",
            "",
            "The JSON report contains slide and layout width values and types, complete spPr XML, complete xfrm XML, layout xfrm XML, and every intermediate mutation state.",
        ]
    )
    md_path.write_text("\n".join(lines), encoding="utf8")

    print(json.dumps({
        "json_report": str(json_path),
        "markdown_report": str(md_path),
        "conclusions": conclusions,
    }, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
