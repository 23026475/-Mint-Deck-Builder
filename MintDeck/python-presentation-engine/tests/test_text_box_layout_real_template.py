"""Real-template regression coverage for inherited placeholder geometry.

This test exercises the production text path:

    approved template
    -> real slide layout
    -> real inherited placeholder
    -> PlaceholderBuilder.populate()
    -> TextBoxLayoutPolicy.apply()
    -> save
    -> reopen

Fake shapes are intentionally not used.
"""

from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

from presentation_engine.builders.placeholder_builder import PlaceholderBuilder


PROJECT_ROOT = Path(__file__).resolve().parents[2]
TEMPLATE_CANDIDATES = (
    PROJECT_ROOT / "data" / "work" / "template" / "FY27 AI-Ready.work.pptx",
    PROJECT_ROOT / "data" / "work" / "template" / "FY27 AI-Ready v3.0.pptx",
)

TARGETS = (
    {
        "layout_name": "Statement – Full Bleed",
        "idx": 2,
        "archetype": "statement",
        "slide_definition": {
            "archetype": "statement",
            "action_title": "IT capability is expanding from infrastructure into data and AI",
            "fields": {
                "statement": (
                    "IT capability is expanding beyond traditional infrastructure into software, data, analytics and AI.\n"
                    "Technology teams now need strong foundations, practical engineering discipline, reliable data practices and automation-ready ways of working.\n"
                    "Takeaway: build a balanced skills portfolio that keeps the core stable while growing intelligent delivery capability."
                ),
                "kicker": "EXECUTIVE OVERVIEW",
            },
        },
        "expected_width_inches": 9.00,
    },
    {
        "layout_name": "Content – Image Right",
        "idx": 12,
        "archetype": "image_right",
        "slide_definition": {
            "archetype": "image_right",
            "action_title": "Technology capability is ultimately enabled by people",
            "fields": {
                "body": (
                    "Capability grows when teams combine engineering discipline, shared learning and practical collaboration. "
                    "This slide validates ImageCatalog resolution, ImagePlaceholderMapper mapping and ImageHandler insertion."
                ),
                "kicker": "IMAGE HANDLER TEST",
            },
        },
        "expected_width_inches": 7.0149,
    },
    {
        "layout_name": "Closing – Next Steps",
        "idx": 12,
        "archetype": "closing",
        "slide_definition": {
            "archetype": "closing",
            "action_title": "Building the capability to move from technology to intelligence",
            "fields": {
                "steps": [
                    "Building technology capability creates the foundation for moving from technology to intelligence.",
                    "Strong foundations across infrastructure, software, data and analytics enable AI-ready delivery.",
                    "Continue expanding the engine's media capabilities.",
                    "Balanced IT skills keep the core stable while growing AI-ready delivery."
                ],
                "kicker": "EXECUTIVE OVERVIEW",
                "footer": "Full engine validation scenario",
            },
        },
        "expected_width_inches": 11.50,
    },
)


class RealTemplateTextBoxGeometryTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls) -> None:
        cls.template_path = next(
            (path for path in TEMPLATE_CANDIDATES if path.exists()),
            None,
        )
        if cls.template_path is None:
            raise unittest.SkipTest(
                "Approved converted template was not found under data/work/template."
            )

    def test_dynamic_placeholders_preserve_effective_geometry_after_save_and_reopen(self):
        presentation = Presentation(self.template_path)
        expected: list[dict[str, int | float]] = []

        for target in TARGETS:
            layout = next(
                (
                    item
                    for item in presentation.slide_layouts
                    if item.name == target["layout_name"]
                ),
                None,
            )
            self.assertIsNotNone(
                layout,
                f"Missing approved layout: {target['layout_name']}",
            )

            slide = presentation.slides.add_slide(layout)
            placeholder = next(
                (
                    item
                    for item in slide.placeholders
                    if int(item.placeholder_format.idx) == target["idx"]
                ),
                None,
            )
            self.assertIsNotNone(
                placeholder,
                f"Missing placeholder idx {target['idx']} on {target['layout_name']}",
            )

            original_left = int(placeholder.left)
            original_top = int(placeholder.top)
            original_width = int(placeholder.width)
            original_height = int(placeholder.height)

            # Authoritative path: this performs text population, template
            # property restoration, text-frame configuration, and the real
            # TextBoxLayoutPolicy integration used by production.
            builder = PlaceholderBuilder()
            builder.populate(
                slide,
                target["slide_definition"],
            )

            layout_entry = next(
                entry
                for entry in builder.last_cleanup_report["populated"]
                if entry["idx"] == target["idx"]
            )
            expected_height_inches = layout_entry["text_box_layout"][
                "applied_height_inches"
            ]
            self.assertGreater(expected_height_inches, 0)

            resized = next(
                item
                for item in slide.placeholders
                if int(item.placeholder_format.idx) == target["idx"]
            )

            self.assertEqual(int(resized.left), original_left)
            self.assertEqual(int(resized.top), original_top)
            self.assertEqual(int(resized.width), original_width)
            self.assertNotEqual(int(resized.height), original_height)
            self.assertAlmostEqual(
                resized.width / 914400,
                target["expected_width_inches"],
                places=2,
            )
            self.assertAlmostEqual(
                resized.height / 914400,
                expected_height_inches,
                places=2,
            )

            xfrm = resized._element.spPr.xfrm
            self.assertIsNotNone(xfrm)
            off = xfrm.find(qn("a:off"))
            ext = xfrm.find(qn("a:ext"))
            self.assertIsNotNone(off)
            self.assertIsNotNone(ext)
            self.assertEqual(int(off.get("x")), original_left)
            self.assertEqual(int(off.get("y")), original_top)
            self.assertEqual(int(ext.get("cx")), original_width)
            self.assertEqual(int(ext.get("cy")), int(resized.height))
            self.assertNotEqual(int(ext.get("cx")), 0)

            expected.append(
                {
                    "slide_index": len(presentation.slides) - 1,
                    "idx": target["idx"],
                    "left": original_left,
                    "top": original_top,
                    "width": original_width,
                    "height": int(resized.height),
                }
            )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = Path(temp_dir) / "real-placeholder-geometry.pptx"
            presentation.save(output_path)
            reopened = Presentation(output_path)

            for item in expected:
                slide = reopened.slides[item["slide_index"]]
                placeholder = next(
                    shape
                    for shape in slide.placeholders
                    if int(shape.placeholder_format.idx) == item["idx"]
                )

                self.assertEqual(int(placeholder.left), item["left"])
                self.assertEqual(int(placeholder.top), item["top"])
                self.assertEqual(int(placeholder.width), item["width"])
                self.assertEqual(int(placeholder.height), item["height"])

                xfrm = placeholder._element.spPr.xfrm
                self.assertIsNotNone(xfrm)
                off = xfrm.find(qn("a:off"))
                ext = xfrm.find(qn("a:ext"))
                self.assertIsNotNone(off)
                self.assertIsNotNone(ext)
                self.assertEqual(int(off.get("x")), item["left"])
                self.assertEqual(int(off.get("y")), item["top"])
                self.assertEqual(int(ext.get("cx")), item["width"])
                self.assertEqual(int(ext.get("cy")), item["height"])
                self.assertNotEqual(int(ext.get("cx")), 0)


if __name__ == "__main__":
    unittest.main()
