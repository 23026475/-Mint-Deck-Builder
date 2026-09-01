from __future__ import annotations

import unittest
from types import SimpleNamespace

from pptx.util import Inches

from presentation_engine.handlers.chart_handler import (
    ChartHandler,
    ChartPlaceholderNotFoundException,
    InvalidChartDataException,
)


class FakeChartTitle:
    def __init__(self) -> None:
        self.text_frame = SimpleNamespace(text="")


class FakeChartObject:
    def __init__(self) -> None:
        self.has_title = False
        self.has_legend = True
        self.chart_title = FakeChartTitle()
        self.series = []


class FakeChartShape:
    def __init__(self, chart_type, left, top, width, height, chart_data) -> None:
        self.name = "Inserted Chart"
        self.chart_type = chart_type
        self.chart_data = chart_data
        self.chart = FakeChartObject()
        self.left = left
        self.top = top
        self.width = width
        self.height = height


class FakeShapes:
    def __init__(self) -> None:
        self.added_charts: list[FakeChartShape] = []

    def add_chart(self, chart_type, left, top, width, height, chart_data) -> FakeChartShape:
        shape = FakeChartShape(chart_type, left, top, width, height, chart_data)
        self.added_charts.append(shape)
        return shape


class FakeElementParent:
    def __init__(self) -> None:
        self.removed_elements: list[object] = []

    def remove(self, element: object) -> None:
        self.removed_elements.append(element)


class FakeElement:
    def __init__(self, parent: FakeElementParent) -> None:
        self.parent = parent

    def getparent(self) -> FakeElementParent:
        return self.parent


class FakeChartPlaceholder:
    def __init__(self, idx: int, width=300, height=200) -> None:
        self.placeholder_format = SimpleNamespace(idx=idx)
        self.left = 10
        self.top = 20
        self.width = width
        self.height = height
        self.parent = FakeElementParent()
        self._element = FakeElement(self.parent)


class FakeSlide:
    def __init__(self, placeholders: list[object]) -> None:
        self.placeholders = placeholders
        self.shapes = FakeShapes()


class ChartHandlerTests(unittest.TestCase):
    def test_insert_column_chart_into_configured_placeholder_using_slide_shapes(self):
        placeholder = FakeChartPlaceholder(idx=10)
        slide = FakeSlide([placeholder])
        handler = ChartHandler()

        report = handler.insert_charts(
            slide,
            {
                "archetype": "chart",
                "fields": {
                    "chart": {
                        "type": "column",
                        "title": "Milestone Progress",
                        "categories": ["Images", "Tables", "Charts"],
                        "series": [
                            {"name": "Status", "values": [100, 60, 10]},
                        ],
                    }
                },
            },
        )

        self.assertEqual(len(report), 1)
        self.assertEqual(report[0]["status"], "inserted")
        self.assertEqual(report[0]["placeholder_idx"], 10)
        self.assertEqual(report[0]["chart_type"], "column")
        self.assertEqual(report[0]["chart_title"], "Milestone Progress")
        self.assertEqual(report[0]["category_count"], 3)
        self.assertEqual(report[0]["series_count"], 1)
        self.assertEqual(report[0]["categories"], ["Images", "Tables", "Charts"])
        self.assertEqual(report[0]["series"][0]["values"], [100.0, 60.0, 10.0])

        self.assertEqual(len(slide.shapes.added_charts), 1)
        inserted = slide.shapes.added_charts[0]
        self.assertEqual(inserted.left, 10)
        self.assertEqual(inserted.top, 20)
        self.assertEqual(inserted.width, handler.MIN_CHART_WIDTH)
        self.assertEqual(inserted.height, handler.MIN_CHART_HEIGHT)
        self.assertTrue(inserted.chart.has_title)
        self.assertEqual(inserted.chart.chart_title.text_frame.text, "Milestone Progress")
        self.assertEqual(placeholder.parent.removed_elements, [placeholder._element])

    def test_sufficient_placeholder_geometry_is_preserved(self):
        placeholder = FakeChartPlaceholder(idx=10, width=Inches(7.0), height=Inches(4.0))
        slide = FakeSlide([placeholder])
        handler = ChartHandler()

        handler.insert_charts(
            slide,
            {
                "archetype": "chart",
                "fields": {
                    "chart": {
                        "type": "column",
                        "categories": ["Images", "Tables", "Charts"],
                        "series": [{"name": "Status", "values": [100, 60, 10]}],
                    }
                },
            },
        )

        inserted = slide.shapes.added_charts[0]
        self.assertEqual(inserted.width, Inches(7.0))
        self.assertEqual(inserted.height, Inches(4.0))

    def test_insert_bar_chart_with_multiple_series(self):
        placeholder = FakeChartPlaceholder(idx=10)
        slide = FakeSlide([placeholder])
        handler = ChartHandler()

        report = handler.insert_charts(
            slide,
            {
                "archetype": "chart",
                "fields": {
                    "chart": {
                        "type": "bar",
                        "categories": ["Images", "Tables", "Charts"],
                        "series": [
                            {"name": "Current", "values": [100, 60, 10]},
                            {"name": "Target", "values": [100, 100, 100]},
                        ],
                    }
                },
            },
        )

        self.assertEqual(report[0]["status"], "inserted")
        self.assertEqual(report[0]["chart_type"], "bar")
        self.assertEqual(report[0]["series_count"], 2)

    def test_missing_chart_payload_produces_no_requests(self):
        placeholder = FakeChartPlaceholder(idx=10)
        slide = FakeSlide([placeholder])
        handler = ChartHandler()

        report = handler.insert_charts(
            slide,
            {
                "archetype": "chart",
                "fields": {},
            },
        )

        self.assertEqual(report, [])
        self.assertEqual(len(slide.shapes.added_charts), 0)

    def test_empty_chart_is_skipped(self):
        placeholder = FakeChartPlaceholder(idx=10)
        slide = FakeSlide([placeholder])
        handler = ChartHandler()

        report = handler.insert_charts(
            slide,
            {
                "archetype": "chart",
                "fields": {
                    "chart": {
                        "type": "column",
                        "categories": [],
                        "series": [],
                    }
                },
            },
        )

        self.assertEqual(len(report), 1)
        self.assertEqual(report[0]["status"], "skipped_empty_chart")
        self.assertEqual(len(slide.shapes.added_charts), 0)
        self.assertEqual(placeholder.parent.removed_elements, [])

    def test_missing_placeholder_raises_clear_error(self):
        slide = FakeSlide([])
        handler = ChartHandler()

        with self.assertRaises(ChartPlaceholderNotFoundException):
            handler.insert_charts(
                slide,
                {
                    "archetype": "chart",
                    "fields": {
                        "chart": {
                            "type": "column",
                            "categories": ["Images"],
                            "series": [{"name": "Status", "values": [100]}],
                        }
                    },
                },
            )

    def test_malformed_chart_payload_raises_clear_error(self):
        placeholder = FakeChartPlaceholder(idx=10)
        slide = FakeSlide([placeholder])
        handler = ChartHandler()

        with self.assertRaises(InvalidChartDataException):
            handler.insert_charts(
                slide,
                {
                    "archetype": "chart",
                    "fields": {"chart": "not a chart object"},
                },
            )

    def test_missing_categories_raises_clear_error_when_series_exists(self):
        placeholder = FakeChartPlaceholder(idx=10)
        slide = FakeSlide([placeholder])
        handler = ChartHandler()

        with self.assertRaises(InvalidChartDataException):
            handler.insert_charts(
                slide,
                {
                    "archetype": "chart",
                    "fields": {
                        "chart": {
                            "type": "column",
                            "series": [{"name": "Status", "values": [100]}],
                        }
                    },
                },
            )

    def test_mismatched_category_and_value_lengths_raise_clear_error(self):
        placeholder = FakeChartPlaceholder(idx=10)
        slide = FakeSlide([placeholder])
        handler = ChartHandler()

        with self.assertRaises(InvalidChartDataException):
            handler.insert_charts(
                slide,
                {
                    "archetype": "chart",
                    "fields": {
                        "chart": {
                            "type": "column",
                            "categories": ["Images", "Tables"],
                            "series": [{"name": "Status", "values": [100]}],
                        }
                    },
                },
            )

    def test_non_numeric_values_raise_clear_error(self):
        placeholder = FakeChartPlaceholder(idx=10)
        slide = FakeSlide([placeholder])
        handler = ChartHandler()

        with self.assertRaises(InvalidChartDataException):
            handler.insert_charts(
                slide,
                {
                    "archetype": "chart",
                    "fields": {
                        "chart": {
                            "type": "column",
                            "categories": ["Images"],
                            "series": [{"name": "Status", "values": ["not numeric"]}],
                        }
                    },
                },
            )


if __name__ == "__main__":
    unittest.main()
