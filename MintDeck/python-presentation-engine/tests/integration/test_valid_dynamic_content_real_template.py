"""Real-template evidence for valid dynamic content after contract enforcement."""

from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

from presentation_engine.builders.placeholder_builder import PlaceholderBuilder


PROJECT_ROOT = Path(__file__).resolve().parents[2]
TEMPLATE_CANDIDATES = (
    PROJECT_ROOT / "data" / "work" / "template" / "FY27 AI-Ready.work.pptx",
    PROJECT_ROOT / "data" / "work" / "template" / "FY27 AI-Ready v3.0.pptx",
)


class ValidDynamicContentRealTemplateTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls) -> None:
        cls.template_path = next(
            (path for path in TEMPLATE_CANDIDATES if path.exists()),
            None,
        )
        if cls.template_path is None:
            raise unittest.SkipTest("Approved converted template not found.")

    def test_valid_statement_support_is_resized_and_preserves_geometry(self):
        presentation = Presentation(self.template_path)
        layout = next(
            item
            for item in presentation.slide_layouts
            if item.name == "Statement – Full Bleed"
        )
        slide = presentation.slides.add_slide(layout)
        placeholder = next(
            item
            for item in slide.placeholders
            if int(item.placeholder_format.idx) == 2
        )

        original = {
            "left": int(placeholder.left),
            "top": int(placeholder.top),
            "width": int(placeholder.width),
            "height": int(placeholder.height),
        }
        valid_text = "Balanced capability strengthens reliable, AI-ready delivery."
        self.assertLessEqual(len(valid_text), 90)
        self.assertEqual(len(valid_text.splitlines()), 1)

        builder = PlaceholderBuilder()
        builder.populate(
            slide,
            {
                "archetype": "statement",
                "action_title": "IT capability is expanding from infrastructure into data and AI",
                "fields": {
                    "statement": valid_text,
                    "kicker": "EXECUTIVE OVERVIEW",
                },
            },
        )

        resized = next(
            item
            for item in slide.placeholders
            if int(item.placeholder_format.idx) == 2
        )
        entry = next(
            item
            for item in builder.last_cleanup_report["populated"]
            if item["idx"] == 2
        )
        layout_result = entry["text_box_layout"]

        self.assertEqual(int(resized.left), original["left"])
        self.assertEqual(int(resized.top), original["top"])
        self.assertEqual(int(resized.width), original["width"])
        self.assertGreater(int(resized.height), 0)
        self.assertAlmostEqual(
            resized.height / 914400,
            layout_result["applied_height_inches"],
            places=3,
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "valid-dynamic-content.pptx"
            presentation.save(output)
            reopened = Presentation(output)
            reopened_slide = reopened.slides[len(reopened.slides) - 1]
            reopened_placeholder = next(
                item
                for item in reopened_slide.placeholders
                if int(item.placeholder_format.idx) == 2
            )

            self.assertEqual(int(reopened_placeholder.left), original["left"])
            self.assertEqual(int(reopened_placeholder.top), original["top"])
            self.assertEqual(int(reopened_placeholder.width), original["width"])
            self.assertEqual(int(reopened_placeholder.height), int(resized.height))

            xfrm = reopened_placeholder._element.spPr.xfrm
            off = xfrm.find(qn("a:off"))
            ext = xfrm.find(qn("a:ext"))
            self.assertEqual(int(off.get("x")), original["left"])
            self.assertEqual(int(off.get("y")), original["top"])
            self.assertEqual(int(ext.get("cx")), original["width"])
            self.assertEqual(int(ext.get("cy")), int(resized.height))
            self.assertNotEqual(int(ext.get("cx")), 0)

        # Evidence for the later safety-margin decision. The test deliberately
        # records the current formula output without imposing a new margin.
        self.assertIn("estimated_height_inches", layout_result)
        self.assertIn("applied_height_inches", layout_result)


if __name__ == "__main__":
    unittest.main()
