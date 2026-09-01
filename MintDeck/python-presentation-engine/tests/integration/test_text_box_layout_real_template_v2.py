"""Real-template regression coverage for inherited placeholder geometry.

This test exercises the production text path using an approved template,
PlaceholderBuilder.populate(), save, and reopen. Fake shapes are not used.
"""

from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from presentation_engine.builders.placeholder_builder import PlaceholderBuilder


PROJECT_ROOT = Path(__file__).resolve().parents[2]
TEMPLATE_CANDIDATES = (
    PROJECT_ROOT / "data" / "work" / "template" / "FY27 AI-Ready.work.pptx",
    PROJECT_ROOT / "data" / "work" / "template" / "FY27 AI-Ready v3.0.pptx",
)

TARGETS = (
    {
        "layout_name": "Statement – Full Bleed",
        "idx": 2,
        "archetype": "statement",
        "slide_definition": {
            "archetype": "statement",
            "action_title": "IT capability is expanding into data and AI",
            "fields": {
                "statement": "Balanced IT skills keep the core stable while growing AI-ready delivery.",
                "kicker": "EXECUTIVE OVERVIEW",
            },
        },
        "expected_width_inches": 9.00,
    },
    {
        "layout_name": "Content – Image Right",
        "idx": 12,
        "archetype": "image_right",
        "slide_definition": {
            "archetype": "image_right",
            "action_title": "Technology capability is enabled by people",
            "fields": {
                "body": "Teams grow capability through disciplined engineering and shared learning.",
                "kicker": "IMAGE HANDLER TEST",
            },
        },
        "expected_width_inches": 7.0149,
    },
    {
        "layout_name": "Closing – Next Steps",
        "idx": 12,
        "archetype": "closing",
        "slide_definition": {
            "archetype": "closing",
            "action_title": "Move from technology to intelligence",
            "fields": {
                "steps": [
                    "Keep the technology foundation stable.",
                    "Grow software, data and analytics skills.",
                    "Continue validating media capabilities.",
                ],
                "kicker": "CLOSING TEST",
                "footer": "Engine validation scenario",
            },
        },
        "expected_width_inches": 11.50,
    },
)


class RealTemplateTextBoxGeometryTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls) -> None:
        cls.template_path = next(
            (path for path in TEMPLATE_CANDIDATES if path.exists()),
            None,
        )
        if cls.template_path is None:
            raise unittest.SkipTest(
                "Approved converted template was not found under data/work/template."
            )

    def test_dynamic_placeholders_preserve_effective_geometry_after_save_and_reopen(self):
        presentation = Presentation(self.template_path)
        expected = []

        for target in TARGETS:
            layout = next(
                (item for item in presentation.slide_layouts if item.name == target["layout_name"]),
                None,
            )
            self.assertIsNotNone(layout, f"Missing approved layout: {target['layout_name']}")

            slide = presentation.slides.add_slide(layout)
            placeholder = next(
                (
                    item
                    for item in slide.placeholders
                    if int(item.placeholder_format.idx) == target["idx"]
                ),
                None,
            )
            self.assertIsNotNone(
                placeholder,
                f"Missing placeholder idx {target['idx']} on {target['layout_name']}",
            )

            builder = PlaceholderBuilder()
            builder.populate(slide, target["slide_definition"])

            populated = next(
                item
                for item in slide.placeholders
                if int(item.placeholder_format.idx) == target["idx"]
            )
            self.assertGreater(int(populated.width), 0)
            self.assertGreater(int(populated.height), 0)

            expected.append(
                {
                    "slide_index": len(presentation.slides) - 1,
                    "idx": target["idx"],
                    "left": int(populated.left),
                    "top": int(populated.top),
                    "width": int(populated.width),
                    "height": int(populated.height),
                }
            )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = Path(temp_dir) / "geometry-regression.pptx"
            presentation.save(output_path)
            reopened = Presentation(output_path)

            for item in expected:
                slide = reopened.slides[item["slide_index"]]
                placeholder = next(
                    shape
                    for shape in slide.placeholders
                    if int(shape.placeholder_format.idx) == item["idx"]
                )
                self.assertEqual(int(placeholder.left), item["left"])
                self.assertEqual(int(placeholder.top), item["top"])
                self.assertEqual(int(placeholder.width), item["width"])
                self.assertEqual(int(placeholder.height), item["height"])


if __name__ == "__main__":
    unittest.main()
