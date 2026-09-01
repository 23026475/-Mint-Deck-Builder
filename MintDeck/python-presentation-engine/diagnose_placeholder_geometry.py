"""Diagnose inherited PowerPoint placeholder geometry through the real engine.

This script does not modify production source files. It instruments the current
runtime in memory, executes the existing DeckBuilder pipeline, and captures the
selected placeholders at these stages:

A  matching layout placeholder used as the inheritance source
B  slide placeholder immediately before TextBoxLayoutPolicy.apply()
C  slide placeholder immediately after TextBoxLayoutPolicy.apply()
C2 presentation state immediately before save()
D  saved placeholder after reopening the generated PPTX

Run from the python-presentation-engine project root:

    $env:PYTHONPATH = ".\\src"
    python .\\diagnose_placeholder_geometry.py

Optional:

    python .\\diagnose_placeholder_geometry.py --contract data/input/it_skills_engine_validation_contract.json
"""

from __future__ import annotations

import argparse
import inspect
import json
import traceback
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Any, Callable

from pptx import Presentation
from pptx.presentation import Presentation as PresentationClass

from presentation_engine.builders.deck_builder import DeckBuilder
from presentation_engine.builders.placeholder_builder import PlaceholderBuilder
from presentation_engine.builders.slide_builder import SlideBuilder
from presentation_engine.services.layout_policy import configure_text_frame
from presentation_engine.services.text_box_layout import TextBoxLayoutPolicy


EMU_PER_INCH = 914400
TARGETS = {
    ("statement", "support"): {"slide_number": 2, "idx": 2},
    ("image_right", "body"): {"slide_number": 6, "idx": 12},
    ("closing", "steps"): {"slide_number": 8, "idx": 12},
}


@dataclass
class GeometrySnapshot:
    stage: str
    slide_number: int
    archetype: str
    placeholder_name: str
    idx: int | None
    shape_name: str | None
    placeholder_type: str | None
    left_emu: int | None
    top_emu: int | None
    width_emu: int | None
    height_emu: int | None
    left_inches: float | None
    top_inches: float | None
    width_inches: float | None
    height_inches: float | None
    has_explicit_xfrm: bool
    xfrm_xml: str | None
    sppr_xml: str | None
    layout_shape_name: str | None
    layout_idx: int | None
    layout_type: str | None
    layout_left_emu: int | None
    layout_top_emu: int | None
    layout_width_emu: int | None
    layout_height_emu: int | None
    layout_left_inches: float | None
    layout_top_inches: float | None
    layout_width_inches: float | None
    layout_height_inches: float | None
    layout_has_explicit_xfrm: bool
    layout_xfrm_xml: str | None
    layout_sppr_xml: str | None


def emu(value: Any) -> int | None:
    if value is None:
        return None
    try:
        return int(value)
    except (TypeError, ValueError):
        return None


def inches(value: Any) -> float | None:
    value_emu = emu(value)
    if value_emu is None:
        return None
    return round(value_emu / EMU_PER_INCH, 4)


def placeholder_idx(shape: Any) -> int | None:
    try:
        return int(shape.placeholder_format.idx)
    except (AttributeError, TypeError, ValueError):
        return None


def placeholder_type(shape: Any) -> str | None:
    try:
        return str(shape.placeholder_format.type)
    except (AttributeError, TypeError, ValueError):
        return None


def xml_for(element: Any) -> str | None:
    if element is None:
        return None
    try:
        return element.xml
    except AttributeError:
        try:
            from lxml import etree
            return etree.tostring(element, pretty_print=True, encoding="unicode")
        except Exception:
            return str(element)


def sppr(shape: Any) -> Any:
    element = getattr(shape, "_element", None)
    if element is None:
        return None
    return getattr(element, "spPr", None)


def xfrm(shape: Any) -> Any:
    shape_sppr = sppr(shape)
    if shape_sppr is None:
        return None
    return getattr(shape_sppr, "xfrm", None)


def slide_from_shape(shape: Any) -> Any | None:
    part = getattr(shape, "part", None)
    return getattr(part, "slide", None)


def matching_layout_placeholder(shape: Any) -> Any | None:
    slide = slide_from_shape(shape)
    if slide is None:
        return None
    idx = placeholder_idx(shape)
    if idx is None:
        return None
    for candidate in slide.slide_layout.placeholders:
        if placeholder_idx(candidate) == idx:
            return candidate
    return None


def snapshot(
    *,
    stage: str,
    shape: Any,
    archetype: str,
    placeholder_name: str,
    slide_number: int,
) -> GeometrySnapshot:
    layout = matching_layout_placeholder(shape)
    shape_xfrm = xfrm(shape)
    layout_xfrm = xfrm(layout) if layout is not None else None

    return GeometrySnapshot(
        stage=stage,
        slide_number=slide_number,
        archetype=archetype,
        placeholder_name=placeholder_name,
        idx=placeholder_idx(shape),
        shape_name=getattr(shape, "name", None),
        placeholder_type=placeholder_type(shape),
        left_emu=emu(getattr(shape, "left", None)),
        top_emu=emu(getattr(shape, "top", None)),
        width_emu=emu(getattr(shape, "width", None)),
        height_emu=emu(getattr(shape, "height", None)),
        left_inches=inches(getattr(shape, "left", None)),
        top_inches=inches(getattr(shape, "top", None)),
        width_inches=inches(getattr(shape, "width", None)),
        height_inches=inches(getattr(shape, "height", None)),
        has_explicit_xfrm=shape_xfrm is not None,
        xfrm_xml=xml_for(shape_xfrm),
        sppr_xml=xml_for(sppr(shape)),
        layout_shape_name=getattr(layout, "name", None) if layout is not None else None,
        layout_idx=placeholder_idx(layout) if layout is not None else None,
        layout_type=placeholder_type(layout) if layout is not None else None,
        layout_left_emu=emu(getattr(layout, "left", None)) if layout is not None else None,
        layout_top_emu=emu(getattr(layout, "top", None)) if layout is not None else None,
        layout_width_emu=emu(getattr(layout, "width", None)) if layout is not None else None,
        layout_height_emu=emu(getattr(layout, "height", None)) if layout is not None else None,
        layout_left_inches=inches(getattr(layout, "left", None)) if layout is not None else None,
        layout_top_inches=inches(getattr(layout, "top", None)) if layout is not None else None,
        layout_width_inches=inches(getattr(layout, "width", None)) if layout is not None else None,
        layout_height_inches=inches(getattr(layout, "height", None)) if layout is not None else None,
        layout_has_explicit_xfrm=layout_xfrm is not None,
        layout_xfrm_xml=xml_for(layout_xfrm),
        layout_sppr_xml=xml_for(sppr(layout)) if layout is not None else None,
    )


def find_placeholder(slide: Any, idx: int) -> Any | None:
    for shape in slide.placeholders:
        if placeholder_idx(shape) == idx:
            return shape
    return None


def extract_output_path(result: Any) -> Path:
    candidate_names = (
        "output_pptx_path",
        "output_path",
        "pptx_path",
        "path",
    )
    for name in candidate_names:
        value = getattr(result, name, None)
        if value:
            return Path(value)
    if isinstance(result, (str, Path)):
        return Path(result)
    if isinstance(result, dict):
        for name in candidate_names:
            if result.get(name):
                return Path(result[name])
    raise RuntimeError(
        "Could not identify the generated PPTX path from DeckBuilder result: "
        f"{result!r}"
    )


def code_location(obj: Callable[..., Any]) -> dict[str, Any]:
    try:
        file_path = inspect.getsourcefile(obj)
        _, line = inspect.getsourcelines(obj)
        return {"file": file_path, "line": line, "qualname": obj.__qualname__}
    except (OSError, TypeError):
        return {"file": None, "line": None, "qualname": getattr(obj, "__qualname__", repr(obj))}


def geometry_row(item: GeometrySnapshot) -> str:
    def f(value: float | None) -> str:
        return "None" if value is None else f"{value:.2f}"
    return (
        f"| {item.stage} | {item.slide_number} | {item.idx} | "
        f"{f(item.width_inches)} | {f(item.height_inches)} | "
        f"{f(item.left_inches)} | {f(item.top_inches)} | "
        f"{item.has_explicit_xfrm} |"
    )


def determine_failure_stage(items: list[GeometrySnapshot]) -> dict[str, str]:
    conclusions: dict[str, str] = {}
    grouped: dict[tuple[int, int | None], dict[str, GeometrySnapshot]] = {}
    for item in items:
        grouped.setdefault((item.slide_number, item.idx), {})[item.stage] = item

    for key, stages in grouped.items():
        label = f"slide_{key[0]}_idx_{key[1]}"
        b = stages.get("B_BEFORE_POLICY")
        c = stages.get("C_AFTER_POLICY")
        pre = stages.get("C2_PRE_SAVE")
        d = stages.get("D_REOPENED")

        if b and (b.width_emu or 0) == 0:
            conclusions[label] = "Width is already zero before TextBoxLayoutPolicy executes."
        elif b and c and (b.width_emu or 0) > 0 and (c.width_emu or 0) == 0:
            conclusions[label] = "Width becomes zero immediately inside TextBoxLayoutPolicy.apply()."
        elif c and pre and (c.width_emu or 0) > 0 and (pre.width_emu or 0) == 0:
            conclusions[label] = "Width becomes zero after policy execution but before presentation save()."
        elif pre and d and (pre.width_emu or 0) > 0 and (d.width_emu or 0) == 0:
            conclusions[label] = "Width becomes zero during save/serialization or reopen."
        elif d and (d.width_emu or 0) > 0:
            conclusions[label] = "Width remains non-zero through reopen; the reported zero came from a different output or inspection path."
        else:
            conclusions[label] = "The captured stages are incomplete; no evidence-based failure stage can be assigned."

    return conclusions


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--contract",
        default="data/input/it_skills_engine_validation_contract.json",
    )
    parser.add_argument(
        "--json-output",
        default="data/output/placeholder_geometry_diagnostic.json",
    )
    parser.add_argument(
        "--markdown-output",
        default="data/output/placeholder_geometry_diagnostic.md",
    )
    args = parser.parse_args()

    contract_path = Path(args.contract)
    if not contract_path.exists():
        raise FileNotFoundError(contract_path)

    snapshots: list[GeometrySnapshot] = []
    runtime_stacks: list[dict[str, Any]] = []

    original_apply = TextBoxLayoutPolicy.apply
    original_save = PresentationClass.save

    def instrumented_apply(
        self: TextBoxLayoutPolicy,
        shape: Any,
        *,
        archetype: str,
        placeholder_name: str,
        text: str,
        font_size_pt: float | None = None,
    ) -> Any:
        target = TARGETS.get((archetype, placeholder_name))
        if target:
            snapshots.append(
                snapshot(
                    stage="A_LAYOUT_SOURCE",
                    shape=matching_layout_placeholder(shape) or shape,
                    archetype=archetype,
                    placeholder_name=placeholder_name,
                    slide_number=target["slide_number"],
                )
            )
            snapshots.append(
                snapshot(
                    stage="B_BEFORE_POLICY",
                    shape=shape,
                    archetype=archetype,
                    placeholder_name=placeholder_name,
                    slide_number=target["slide_number"],
                )
            )
            runtime_stacks.append(
                {
                    "slide_number": target["slide_number"],
                    "idx": target["idx"],
                    "stack": traceback.format_stack(limit=12),
                }
            )

        result = original_apply(
            self,
            shape,
            archetype=archetype,
            placeholder_name=placeholder_name,
            text=text,
            font_size_pt=font_size_pt,
        )

        if target:
            snapshots.append(
                snapshot(
                    stage="C_AFTER_POLICY",
                    shape=shape,
                    archetype=archetype,
                    placeholder_name=placeholder_name,
                    slide_number=target["slide_number"],
                )
            )
        return result

    def instrumented_save(self: Any, file: Any) -> Any:
        for (archetype, placeholder_name), target in TARGETS.items():
            if target["slide_number"] > len(self.slides):
                continue
            slide = self.slides[target["slide_number"] - 1]
            shape = find_placeholder(slide, target["idx"])
            if shape is not None:
                snapshots.append(
                    snapshot(
                        stage="C2_PRE_SAVE",
                        shape=shape,
                        archetype=archetype,
                        placeholder_name=placeholder_name,
                        slide_number=target["slide_number"],
                    )
                )
        return original_save(self, file)

    TextBoxLayoutPolicy.apply = instrumented_apply
    PresentationClass.save = instrumented_save

    try:
        result = DeckBuilder().build_from_contract_file(contract_path)
        output_path = extract_output_path(result)
    finally:
        TextBoxLayoutPolicy.apply = original_apply
        PresentationClass.save = original_save

    reopened = Presentation(output_path)
    for (archetype, placeholder_name), target in TARGETS.items():
        slide = reopened.slides[target["slide_number"] - 1]
        shape = find_placeholder(slide, target["idx"])
        if shape is not None:
            snapshots.append(
                snapshot(
                    stage="D_REOPENED",
                    shape=shape,
                    archetype=archetype,
                    placeholder_name=placeholder_name,
                    slide_number=target["slide_number"],
                )
            )

    architecture = {
        "template_to_output_code_locations": [
            code_location(DeckBuilder.build_from_contract_file),
            code_location(DeckBuilder.build),
            code_location(SlideBuilder.build_slide),
            code_location(PlaceholderBuilder.populate),
            code_location(PlaceholderBuilder._populate_from_definitions),
            code_location(PlaceholderBuilder._placeholder_by_idx),
            code_location(PlaceholderBuilder._fill_text_placeholder),
            code_location(configure_text_frame),
            code_location(TextBoxLayoutPolicy.apply),
            code_location(PresentationClass.save),
        ],
        "runtime_stacks": runtime_stacks,
    }

    conclusions = determine_failure_stage(snapshots)
    payload = {
        "contract": str(contract_path),
        "generated_pptx": str(output_path),
        "snapshots": [asdict(item) for item in snapshots],
        "failure_stage_conclusions": conclusions,
        "architecture": architecture,
    }

    json_output = Path(args.json_output)
    markdown_output = Path(args.markdown_output)
    json_output.parent.mkdir(parents=True, exist_ok=True)
    markdown_output.parent.mkdir(parents=True, exist_ok=True)
    json_output.write_text(json.dumps(payload, indent=2), encoding="utf8")

    lines = [
        "# Placeholder Geometry Diagnostic",
        "",
        f"Generated PPTX: `{output_path}`",
        "",
        "| Stage | Slide | idx | Width (in) | Height (in) | Left (in) | Top (in) | Explicit xfrm |",
        "|---|---:|---:|---:|---:|---:|---:|---|",
    ]
    lines.extend(geometry_row(item) for item in snapshots)
    lines.extend(["", "## Evidence-based failure-stage conclusions", ""])
    for key, value in conclusions.items():
        lines.append(f"- **{key}**: {value}")
    lines.extend(
        [
            "",
            "## OOXML",
            "",
            "The complete slide-placeholder and matching layout-placeholder spPr/xfrm XML is stored in the JSON report.",
        ]
    )
    markdown_output.write_text("\n".join(lines), encoding="utf8")

    print(json.dumps({
        "generated_pptx": str(output_path),
        "json_report": str(json_output),
        "markdown_report": str(markdown_output),
        "failure_stage_conclusions": conclusions,
    }, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
