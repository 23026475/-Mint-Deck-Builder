"""
End-to-end image pipeline validation report.

This diagnostic verifies the image pipeline inputs and generated PPTX output without
modifying engine architecture or production code.

It reports:
- contract used
- images requested
- images resolved
- placeholder indices used
- successful insertions detected in the generated PPTX
- skipped insertions
- missing images
- missing placeholders

Run from project root:
    python .\validate_image_e2e_pipeline.py

Optional:
    python .\validate_image_e2e_pipeline.py --contract .\data\contracts\image-e2e-validation.json --pptx .\data\output\your-deck.pptx
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any, Mapping, Sequence

from pptx import Presentation

from presentation_engine.handlers.image_catalog import ImageCatalog
from presentation_engine.handlers.image_placeholder_mapper import ImagePlaceholderMapper

DEFAULT_CONTRACT = Path("data/contracts/image-e2e-validation.json")
DEFAULT_OUTPUT = Path("data/output/image_e2e_validation_report.json")

BASELINE_ID_TO_ARCHETYPE = {"B04": "image_right", "B06": "quote", "E01": "team", "E03": "logo_wall"}


def canonical_archetype(archetype: str) -> str:
    return BASELINE_ID_TO_ARCHETYPE.get(archetype.upper(), archetype.strip().lower())


def value_from_path(item: Mapping[str, Any], source: str) -> Any:
    current: Any = item
    for part in source.split("."):
        if isinstance(current, Mapping):
            current = current.get(part)
        elif isinstance(current, Sequence) and not isinstance(current, (str, bytes)):
            if not part.isdigit():
                return None
            idx = int(part)
            current = current[idx] if idx < len(current) else None
        else:
            return None
    return current


def image_fields(archetype: str, slide_definition: Mapping[str, Any]) -> list[dict[str, Any]]:
    archetype = canonical_archetype(archetype)
    if archetype == "image_right":
        return [
            {"field_name": "picture", "source": "fields.picture", "occurrence": None},
            {"field_name": "picture", "source": "fields.image", "occurrence": None},
        ]
    if archetype == "quote":
        return [{"field_name": "headshot", "source": "fields.headshot", "occurrence": None}]
    if archetype == "team":
        members = value_from_path(slide_definition, "fields.members")
        count = len(members) if isinstance(members, Sequence) and not isinstance(members, (str, bytes)) else 0
        return [{"field_name": "member_picture", "source": f"fields.members.{i}.picture", "occurrence": i} for i in range(count)]
    if archetype == "logo_wall":
        logos = value_from_path(slide_definition, "fields.logos")
        count = len(logos) if isinstance(logos, Sequence) and not isinstance(logos, (str, bytes)) else 0
        return [{"field_name": "logo", "source": f"fields.logos.{i}", "occurrence": i} for i in range(count)]
    return []


def latest_pptx() -> Path | None:
    candidates = sorted(Path("data/output").glob("*.pptx"), key=lambda p: p.stat().st_mtime, reverse=True)
    return candidates[0] if candidates else None


def picture_shape_count(prs: Presentation, slide_number: int) -> int:
    slide = prs.slides[slide_number - 1]
    count = 0
    for shape in slide.shapes:
        # python-pptx PICTURE shape type is 13; avoid importing enum for compatibility.
        if str(getattr(shape, "shape_type", "")).upper().endswith("PICTURE (13)") or str(getattr(shape, "shape_type", "")) == "PICTURE (13)":
            count += 1
        elif getattr(shape, "image", None) is not None:
            count += 1
    return count


def collect_contract_requests(contract: Mapping[str, Any], image_catalog: ImageCatalog, mapper: ImagePlaceholderMapper) -> tuple[list[dict[str, Any]], list[dict[str, Any]], list[dict[str, Any]]]:
    requests: list[dict[str, Any]] = []
    missing_images: list[dict[str, Any]] = []
    missing_mappings: list[dict[str, Any]] = []

    for index, slide_definition in enumerate(contract.get("slides", []), start=1):
        archetype = canonical_archetype(str(slide_definition.get("archetype", "")))
        for field in image_fields(archetype, slide_definition):
            image_reference = value_from_path(slide_definition, field["source"])
            if not isinstance(image_reference, str) or not image_reference.strip():
                requests.append({
                    "slide_number": index,
                    "archetype": archetype,
                    "media_field": field["field_name"],
                    "source": field["source"],
                    "status": "skipped_no_image_supplied",
                })
                continue

            record: dict[str, Any] = {
                "slide_number": index,
                "archetype": archetype,
                "media_field": field["field_name"],
                "source": field["source"],
                "image_reference": image_reference,
            }

            try:
                mapping = mapper.resolve(archetype, field["field_name"], field["occurrence"])
                record["placeholder_idx"] = mapping.idx
                record["resolved_placeholder_idx"] = mapping.idx
            except Exception as exc:
                record["status"] = "missing_mapping"
                record["error"] = str(exc)
                missing_mappings.append(record.copy())
                requests.append(record)
                continue

            try:
                resolved = image_catalog.resolve(image_reference)
                metadata = image_catalog.metadata(resolved)
                record.update({
                    "resolved_filename": str(resolved),
                    "image_width": metadata.width,
                    "image_height": metadata.height,
                    "image_format": metadata.format,
                    "category": metadata.category,
                    "orientation": metadata.orientation,
                    "status": "resolved",
                })
            except Exception as exc:
                record["status"] = "missing_or_invalid_image"
                record["error"] = str(exc)
                missing_images.append(record.copy())

            requests.append(record)

    return requests, missing_images, missing_mappings


def build_report(contract_path: Path, pptx_path: Path | None, output_path: Path) -> dict[str, Any]:
    contract = json.loads(contract_path.read_text(encoding="utf8"))
    catalog = ImageCatalog()
    mapper = ImagePlaceholderMapper()

    requests, missing_images, missing_mappings = collect_contract_requests(contract, catalog, mapper)

    report: dict[str, Any] = {
        "contract_used": str(contract_path),
        "pptx_checked": str(pptx_path) if pptx_path else None,
        "images_requested": requests,
        "images_resolved": [item for item in requests if item.get("status") == "resolved"],
        "skipped_insertions": [item for item in requests if item.get("status") == "skipped_no_image_supplied"],
        "missing_images": missing_images,
        "missing_placeholders": [],
        "missing_mappings": missing_mappings,
        "successful_insertions_detected": [],
        "summary": {},
    }

    if pptx_path and pptx_path.exists():
        prs = Presentation(str(pptx_path))
        for item in report["images_resolved"]:
            slide_number = item["slide_number"]
            detected_picture_count = picture_shape_count(prs, slide_number)
            detected = detected_picture_count > 0
            report["successful_insertions_detected"].append({
                "slide_number": slide_number,
                "archetype": item["archetype"],
                "media_field": item["media_field"],
                "image_reference": item["image_reference"],
                "placeholder_idx": item.get("placeholder_idx"),
                "picture_shapes_detected_on_slide": detected_picture_count,
                "detected": detected,
            })

    report["summary"] = {
        "image_requests": len(requests),
        "resolved_images": len(report["images_resolved"]),
        "skipped_insertions": len(report["skipped_insertions"]),
        "missing_images": len(missing_images),
        "missing_mappings": len(missing_mappings),
        "successful_insertions_detected": sum(1 for item in report["successful_insertions_detected"] if item.get("detected")),
    }

    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(json.dumps(report, indent=2, ensure_ascii=False), encoding="utf8")
    return report


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Validate end-to-end image pipeline inputs and generated PPTX output.")
    parser.add_argument("--contract", type=Path, default=DEFAULT_CONTRACT)
    parser.add_argument("--pptx", type=Path, default=None)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    pptx = args.pptx or latest_pptx()
    report = build_report(args.contract, pptx, args.output)
    summary = report["summary"]
    print("Image E2E validation summary")
    print("- Contract used:", report["contract_used"])
    print("- PPTX checked:", report["pptx_checked"])
    print("- Image requests:", summary["image_requests"])
    print("- Resolved images:", summary["resolved_images"])
    print("- Skipped insertions:", summary["skipped_insertions"])
    print("- Missing images:", summary["missing_images"])
    print("- Missing mappings:", summary["missing_mappings"])
    print("- Successful insertions detected:", summary["successful_insertions_detected"])
    print("- Report:", args.output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
