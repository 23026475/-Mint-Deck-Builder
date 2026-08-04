
"""
End-to-end validation runner for Mint Presentation Engine exact-match archetypes.
Run from the project root: python-presentation-engine

This script writes outputs to data/output and includes PlaceholderBuilder cleanup reports.
"""
from __future__ import annotations

import json
import re
import zipfile
from collections import Counter
from pathlib import Path

from pptx import Presentation

from presentation_engine.builders.deck_builder import DeckBuilder
from presentation_engine.services.layout_mapper import LayoutMapper
from presentation_engine.services.archetype_mapper import ArchetypeMapper

PROMPT_PATTERNS = [
    re.compile(r"click to add", re.I),
    re.compile(r"click to edit", re.I),
    re.compile(r"click to edit master", re.I),
    re.compile(r"presentation title", re.I),
    re.compile(r"one-line subtitle", re.I),
    re.compile(r"prepared for client name", re.I),
    re.compile(r"section title goes here", re.I),
    re.compile(r"lorem ipsum", re.I),
]

EXPECTED_LAYOUTS = {
    "cover": "Cover – Brand",
    "cover_dark": "Cover – Dark",
    "statement": "Statement – Full Bleed",
    "thesis": "Thesis",
    "cards3": "Content – Cards 3",
    "image_right": "Content – Image Right",
    "comparison": "Content – Comparison",
    "faq": "FAQ",
    "process_flow": "Content – Process Flow",
    "table": "Content – Table",
    "chart": "Content – Chart",
    "kpi": "KPI – 4 Stats",
    "matrix": "Matrix 2x2",
    "team": "Meet the Team",
    "org_chart": "Org Chart",
    "logo_wall": "Logo Wall",
    "case_study": "Case Study",
    "quote": "Quote",
    "pricing_stat": "Content – Pricing Stat",
    "summary_cta": "Content Slide 1 – Body + CTA",
    "closing": "Closing – Next Steps",
}


def duplicate_zip_entries(path: Path) -> dict[str, int]:
    with zipfile.ZipFile(path, "r") as archive:
        counts = Counter(archive.namelist())
    return {name: count for name, count in counts.items() if count > 1}


def shape_texts(slide):
    texts = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = (shape.text or "").strip()
            if text:
                texts.append(text)
    return texts


def has_prompt_leak(texts: list[str]) -> list[str]:
    leaked = []
    for text in texts:
        folded = " ".join(text.split())
        if any(pattern.search(folded) for pattern in PROMPT_PATTERNS):
            leaked.append(text)
    return leaked


def main() -> int:
    contract_path = Path("data/input/exact_match_e2e_contract.json")
    if not contract_path.exists():
        raise SystemExit(f"Missing contract: {contract_path}")

    contract = json.loads(contract_path.read_text(encoding="utf8"))
    builder = DeckBuilder()
    result = builder.build_from_contract_file(contract_path)
    output_path = Path(result.output_pptx_path)

    duplicates = duplicate_zip_entries(output_path)
    prs = Presentation(str(output_path))

    layout_mapper = LayoutMapper()
    layout_mapper.load_from_presentation(prs, output_path)
    archetype_mapper = ArchetypeMapper(layout_mapper=layout_mapper)

    cleanup_reports = getattr(builder.placeholder_builder, "cleanup_reports", [])
    slide_reports = []
    defects = []

    if len(prs.slides) != len(contract["slides"]):
        defects.append(f"Slide count mismatch: expected {len(contract['slides'])}, found {len(prs.slides)}")

    for index, slide_definition in enumerate(contract["slides"]):
        archetype = slide_definition["archetype"]
        slide = prs.slides[index]
        actual_layout = slide.slide_layout.name
        expected_layout = EXPECTED_LAYOUTS.get(archetype)

        try:
            resolved_layout = archetype_mapper.get_layout_name_for_archetype(archetype)
        except Exception as exc:
            resolved_layout = None
            defects.append(f"Slide {index+1} {archetype}: layout resolution failed: {exc}")

        texts = shape_texts(slide)
        prompt_leaks = has_prompt_leak(texts)
        if prompt_leaks:
            defects.append(f"Slide {index+1} {archetype}: unresolved default PowerPoint prompt text found")

        if expected_layout and actual_layout != expected_layout:
            defects.append(f"Slide {index+1} {archetype}: expected layout {expected_layout}, found {actual_layout}")

        if resolved_layout and expected_layout and resolved_layout != expected_layout:
            defects.append(f"Slide {index+1} {archetype}: mapper resolved {resolved_layout}, expected {expected_layout}")

        cleanup_report = cleanup_reports[index] if index < len(cleanup_reports) else {}
        slide_reports.append({
            "slide_number": index + 1,
            "archetype": archetype,
            "expected_layout": expected_layout,
            "actual_layout": actual_layout,
            "mapper_resolved_layout": resolved_layout,
            "text_shape_count": len(texts),
            "has_prompt_leak": bool(prompt_leaks),
            "prompt_leaks": prompt_leaks,
            "populated_placeholders": cleanup_report.get("populated", []),
            "removed_unfilled_text_placeholders": cleanup_report.get("removed_unfilled_text", []),
            "removed_prompt_text_placeholders": cleanup_report.get("removed_prompt_text", []),
            "preserved_media_placeholders": cleanup_report.get("preserved_media", []),
            "status": "PASS" if not prompt_leaks and (not expected_layout or actual_layout == expected_layout) else "FAIL",
        })

    report = {
        "contract": str(contract_path),
        "generated_pptx": str(output_path),
        "slide_count": len(prs.slides),
        "expected_slide_count": len(contract["slides"]),
        "duplicate_zip_entries": len(duplicates),
        "duplicate_zip_entry_details": duplicates,
        "defects": defects,
        "slides": slide_reports,
        "cleanup_report_count": len(cleanup_reports),
        "overall_status": "PASS" if not defects and not duplicates else "FAIL",
    }

    report_path = Path("data/output/exact_match_e2e_validation_report.json")
    report_path.write_text(json.dumps(report, indent=2, ensure_ascii=False), encoding="utf8")

    print(json.dumps(report, indent=2, ensure_ascii=False))
    return 0 if report["overall_status"] == "PASS" else 1


if __name__ == "__main__":
    raise SystemExit(main())
