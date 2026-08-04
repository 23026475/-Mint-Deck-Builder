
"""
Slide builder for the Python Presentation Engine.

Responsibility:
- Resolve layouts through ArchetypeMapper.
- Create slides with python-pptx.
- Delegate placeholder population to PlaceholderBuilder.
"""

from __future__ import annotations

import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Mapping, Optional, Protocol, Sequence

from presentation_engine.config import EngineConfig, config
from presentation_engine.builders.placeholder_builder import (
    PlaceholderBuilder,
    PlaceholderBuilderException,
    SUPPORTED_ARCHETYPES,
)


logger = logging.getLogger(__name__)


class SlideBuilderException(Exception):
    """Raised when slides cannot be created from the supplied contract."""


class ArchetypeLayoutResolver(Protocol):
    def get_layout_for_archetype(self, archetype: str) -> Any:
        """Return the layout object configured for the supplied archetype."""


class PlaceholderPopulator(Protocol):
    def populate(self, slide: Any, slide_definition: Mapping[str, Any]) -> None:
        """Populate placeholders on the supplied slide."""


class SlideCollection(Protocol):
    def add_slide(self, slide_layout: Any) -> Any:
        """Create and return a slide using the supplied layout."""


class PresentationLike(Protocol):
    slides: SlideCollection


class PresentationFactory(Protocol):
    def create_from_template(self, template_path: Path) -> PresentationLike:
        """Load and return a presentation from the converted template."""


@dataclass(frozen=True)
class PythonPptxPresentationFactory:
    def create_from_template(self, template_path: Path) -> PresentationLike:
        try:
            from pptx import Presentation
            return Presentation(str(template_path))
        except Exception as exc:
            raise SlideBuilderException(
                f"Failed to create presentation from converted template: {template_path}"
            ) from exc


@dataclass(frozen=True)
class SlideBuildResult:
    presentation: PresentationLike
    slides: tuple[Any, ...]


@dataclass
class SlideBuilder:
    archetype_mapper: ArchetypeLayoutResolver
    engine_config: EngineConfig = field(default_factory=lambda: config)
    presentation_factory: PresentationFactory = field(default_factory=PythonPptxPresentationFactory)
    placeholder_builder: PlaceholderPopulator = field(default_factory=PlaceholderBuilder)
    supported_archetypes: frozenset[str] = SUPPORTED_ARCHETYPES
    log: logging.Logger = field(default_factory=lambda: logger)

    def build(self, contract: Mapping[str, Any], template_path: Optional[str | Path] = None) -> SlideBuildResult:
        resolved_template_path = self._resolve_template_path(template_path)
        slide_definitions = self._extract_slide_definitions(contract)

        presentation = self.presentation_factory.create_from_template(resolved_template_path)
        created_slides: list[Any] = []
        for slide_definition in slide_definitions:
            created_slides.append(self.build_slide(presentation, slide_definition))
        return SlideBuildResult(presentation=presentation, slides=tuple(created_slides))

    def build_slide(self, presentation: PresentationLike, slide_definition: Mapping[str, Any]) -> Any:
        archetype = self._read_archetype(slide_definition)
        self._validate_supported_archetype(archetype)

        try:
            layout = self.archetype_mapper.get_layout_for_archetype(archetype)
        except Exception as exc:
            raise SlideBuilderException(f"Failed to resolve layout for archetype '{archetype}'.") from exc

        slide = presentation.slides.add_slide(layout)

        try:
            self.placeholder_builder.populate(slide, slide_definition)
        except PlaceholderBuilderException as exc:
            raise SlideBuilderException(f"Failed to populate placeholders for archetype '{archetype}'.") from exc

        return slide

    def _resolve_template_path(self, template_path: Optional[str | Path]) -> Path:
        resolved_path = (
            Path(template_path).expanduser().resolve()
            if template_path is not None
            else self.engine_config.template.converted_pptx_path
        )
        if not resolved_path.exists():
            raise SlideBuilderException(f"Converted FY27 template was not found: {resolved_path}")
        if not resolved_path.is_file():
            raise SlideBuilderException(f"Converted FY27 template path is not a file: {resolved_path}")
        if resolved_path.suffix.lower() != ".pptx":
            raise SlideBuilderException(f"Converted FY27 template must be a .pptx file: {resolved_path}")
        return resolved_path

    def _extract_slide_definitions(self, contract: Mapping[str, Any]) -> Sequence[Mapping[str, Any]]:
        slides = contract.get("slides")
        if not isinstance(slides, list):
            raise SlideBuilderException("JSON contract must contain a 'slides' list.")
        for index, slide_definition in enumerate(slides):
            if not isinstance(slide_definition, Mapping):
                raise SlideBuilderException(f"Slide definition at index {index} must be an object.")
        return slides

    def _read_archetype(self, slide_definition: Mapping[str, Any]) -> str:
        archetype = slide_definition.get("archetype")
        if not isinstance(archetype, str) or not archetype.strip():
            raise SlideBuilderException("Each slide definition must contain a non-empty 'archetype'.")
        return archetype.strip().lower()

    def _validate_supported_archetype(self, archetype: str) -> None:
        if archetype not in self.supported_archetypes:
            supported = sorted(self.supported_archetypes)
            raise SlideBuilderException(f"Unsupported archetype '{archetype}'. Supported archetypes: {supported}")
