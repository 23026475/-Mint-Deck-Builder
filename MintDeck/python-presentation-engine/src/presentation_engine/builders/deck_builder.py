"""
Deck builder orchestration layer for the Python Presentation Engine.

Milestone scope:
- Build the first executable end-to-end deck pipeline.
- Use the local FY27 AI-Ready v3.0 POTX template.
- Convert POTX to PPTX when required.
- Remove sample slides from the converted template.
- Create only the supported milestone slides: cover, cards3, closing.
- Save the generated PPTX to data/output.

Important:
- This module does not populate placeholders.
- This module does not perform QA validation.
- This module does not implement density validation.
- This module does not hardcode layouts or use layout indexes.
"""

from __future__ import annotations

import json
import logging
from dataclasses import dataclass, field
from datetime import date
from pathlib import Path
from typing import Any, Mapping, Optional, Protocol

from presentation_engine.config import EngineConfig, config
from presentation_engine.builders.slide_builder import SlideBuilder, SlideBuilderException
from presentation_engine.services.archetype_mapper import ArchetypeMapper, ArchetypeMappingException
from presentation_engine.services.layout_mapper import LayoutMapper, TemplateLoadError
from presentation_engine.services.potx_converter import PotxConverter, PotxConversionError
from presentation_engine.services.template_retrieval_service import (
    TemplateRetrievalError,
    TemplateRetrievalService,
)


logger = logging.getLogger(__name__)


class DeckBuilderException(Exception):
    """Raised when the deck orchestration pipeline fails."""


class TemplateRetriever(Protocol):
    """Interface for retrieving the approved POTX template."""

    def retrieve_latest_template(self) -> Path:
        """Return the local path to the approved POTX template."""


class TemplateConverter(Protocol):
    """Interface for converting a POTX template to a working PPTX."""

    def convert(self, potx_path: str | Path, pptx_path: str | Path, overwrite: bool = True) -> Path:
        """Convert potx_path to pptx_path and return the PPTX path."""


class PresentationLoader(Protocol):
    """Interface for loading a PPTX presentation."""

    def load(self, presentation_path: Path) -> Any:
        """Load and return a presentation object."""


class SampleSlideRemover(Protocol):
    """Interface for removing all existing sample slides from a presentation."""

    def remove_all(self, presentation: Any) -> None:
        """Remove every slide from the presentation."""


class ContractLoader(Protocol):
    """Interface for loading a parsed JSON contract."""

    def load(self, contract_path: Path) -> Mapping[str, Any]:
        """Load and return a parsed JSON contract."""


class OutputNameProvider(Protocol):
    """Interface for creating output PPTX file names."""

    def get_output_path(self, contract: Mapping[str, Any], output_dir: Path) -> Path:
        """Return the full output PPTX path for the generated deck."""


@dataclass(frozen=True)
class PythonPptxPresentationLoader:
    """Loads a PPTX presentation using python-pptx."""

    def load(self, presentation_path: Path) -> Any:
        """Load a presentation from a PPTX file path."""

        try:
            from pptx import Presentation

            return Presentation(str(presentation_path))
        except Exception as exc:
            raise DeckBuilderException(f"Failed to load converted presentation: {presentation_path}") from exc


@dataclass(frozen=True)
class PythonPptxSampleSlideRemover:
    """
    Removes all sample slides from a python-pptx presentation.

    This follows the safe OpenXML relationship removal sequence used for the
    presentation engine build guide: drop the relationship, then remove the slide
    id entry from the slide id list.
    """

    def remove_all(self, presentation: Any) -> None:
        """Remove every existing slide from a python-pptx presentation."""

        xml_slides = presentation.slides._sldIdLst
        for slide_id in list(xml_slides):
            relationship_id = slide_id.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
            )
            presentation.part.drop_rel(relationship_id)
            xml_slides.remove(slide_id)


@dataclass(frozen=True)
class JsonContractLoader:
    """Loads the Mint Deck Builder JSON contract from disk."""

    def load(self, contract_path: Path) -> Mapping[str, Any]:
        """Load and parse the JSON contract file."""

        if not contract_path.exists():
            raise DeckBuilderException(f"JSON contract file was not found: {contract_path}")

        if not contract_path.is_file():
            raise DeckBuilderException(f"JSON contract path is not a file: {contract_path}")

        try:
            data = json.loads(contract_path.read_text(encoding="utf8"))
        except json.JSONDecodeError as exc:
            raise DeckBuilderException(f"Invalid JSON contract file: {contract_path}") from exc
        except Exception as exc:
            raise DeckBuilderException(f"Failed to read JSON contract file: {contract_path}") from exc

        if not isinstance(data, Mapping):
            raise DeckBuilderException("JSON contract root must be an object.")

        return data


@dataclass(frozen=True)
class DefaultOutputNameProvider:
    """Creates a simple generated deck filename from the JSON contract metadata."""

    fallback_file_name: str = "mint-presentation-engine-milestone.pptx"

    def get_output_path(self, contract: Mapping[str, Any], output_dir: Path) -> Path:
        """Return a safe output PPTX path under data/output."""

        deck = contract.get("deck", {})
        if not isinstance(deck, Mapping):
            return output_dir / self.fallback_file_name

        client = self._safe_name(deck.get("client"))
        title = self._safe_name(deck.get("title"))

        if not client or not title:
            return output_dir / self.fallback_file_name

        return output_dir / f"{client} - {title} - {date.today().isoformat()}.pptx"

    def _safe_name(self, value: Any) -> str:
        if not isinstance(value, str):
            return ""

        unsafe_characters = '<>:"/\\|?*'
        cleaned = value.strip()
        for character in unsafe_characters:
            cleaned = cleaned.replace(character, "")
        return " ".join(cleaned.split())


@dataclass(frozen=True)
class DeckBuildResult:
    """Result returned after the first milestone deck is generated."""

    potx_path: Path
    converted_pptx_path: Path
    output_pptx_path: Path
    slide_count: int


@dataclass
class DeckBuilder:
    """
    Orchestrates the first milestone deck generation pipeline.

    The builder coordinates existing services without duplicating their logic:
    TemplateRetrievalService, PotxConverter, LayoutMapper, ArchetypeMapper, and
    SlideBuilder. It deletes template sample slides, creates milestone slides,
    and saves the output PPTX.
    """

    engine_config: EngineConfig = field(default_factory=lambda: config)
    template_retriever: TemplateRetriever = field(default_factory=TemplateRetrievalService)
    template_converter: TemplateConverter = field(default_factory=PotxConverter)
    presentation_loader: PresentationLoader = field(default_factory=PythonPptxPresentationLoader)
    sample_slide_remover: SampleSlideRemover = field(default_factory=PythonPptxSampleSlideRemover)
    contract_loader: ContractLoader = field(default_factory=JsonContractLoader)
    output_name_provider: OutputNameProvider = field(default_factory=DefaultOutputNameProvider)
    log: logging.Logger = field(default_factory=lambda: logger)

    def build_from_contract_file(self, contract_path: Optional[str | Path] = None) -> DeckBuildResult:
        """
        Build the first milestone presentation from a JSON contract file.

        Args:
            contract_path: Optional explicit path to a JSON contract. When omitted,
                config.input.contract_path is used.

        Returns:
            DeckBuildResult containing key paths and generated slide count.

        Raises:
            DeckBuilderException: If any orchestration step fails.
        """

        resolved_contract_path = self._resolve_contract_path(contract_path)
        self.log.info("Loading JSON contract: %s", resolved_contract_path)
        contract = self.contract_loader.load(resolved_contract_path)
        return self.build(contract)

    def build(self, contract: Mapping[str, Any]) -> DeckBuildResult:
        """
        Build the first milestone presentation from an already parsed contract.

        Args:
            contract: Parsed Mint Deck Builder JSON contract.

        Returns:
            DeckBuildResult with output metadata.
        """

        try:
            self.engine_config.ensure_runtime_directories()
            self.engine_config.output.output_dir.mkdir(parents=True, exist_ok=True)

            self.log.info("Retrieving local FY27 AI-Ready POTX template.")
            potx_path = self.template_retriever.retrieve_latest_template()

            self.log.info("Ensuring converted PPTX exists.")
            converted_pptx_path = self._ensure_converted_template(potx_path)

            self.log.info("Loading layout mapper for converted template.")
            layout_mapper = LayoutMapper(engine_config=self.engine_config)
            layout_mapper.load(converted_pptx_path)

            self.log.info("Loading archetype mapper from baseline configuration.")
            archetype_mapper = ArchetypeMapper(layout_mapper=layout_mapper, engine_config=self.engine_config)

            self.log.info("Loading converted presentation for deck generation.")
            presentation = self.presentation_loader.load(converted_pptx_path)

            self.log.info("Removing all sample slides from template.")
            self.sample_slide_remover.remove_all(presentation)

            self.log.info("Building milestone slides.")
            slide_builder = SlideBuilder(archetype_mapper=archetype_mapper, engine_config=self.engine_config)
            created_slides = []
            for slide_definition in self._get_slide_definitions(contract):
                created_slides.append(slide_builder.build_slide(presentation, slide_definition))

            output_path = self.output_name_provider.get_output_path(
                contract,
                self.engine_config.output.output_dir,
            )
            output_path.parent.mkdir(parents=True, exist_ok=True)

            self.log.info("Saving generated presentation: %s", output_path)
            presentation.save(str(output_path))

            self.log.info("Deck generation completed with %s slides.", len(created_slides))
            return DeckBuildResult(
                potx_path=potx_path,
                converted_pptx_path=converted_pptx_path,
                output_pptx_path=output_path,
                slide_count=len(created_slides),
            )
        except (
            TemplateRetrievalError,
            PotxConversionError,
            TemplateLoadError,
            ArchetypeMappingException,
            SlideBuilderException,
            DeckBuilderException,
        ):
            self.log.exception("Deck generation failed.")
            raise
        except Exception as exc:
            self.log.exception("Unexpected deck generation failure.")
            raise DeckBuilderException("Unexpected failure while generating the milestone deck.") from exc

    def _ensure_converted_template(self, potx_path: Path) -> Path:
        """Convert POTX to PPTX when the converted file is missing or stale."""

        converted_pptx_path = self.engine_config.template.converted_pptx_path

        conversion_required = (
            not converted_pptx_path.exists()
            or potx_path.stat().st_mtime > converted_pptx_path.stat().st_mtime
        )

        if not conversion_required:
            self.log.info("Using existing converted PPTX: %s", converted_pptx_path)
            return converted_pptx_path

        self.log.info("Converting POTX to PPTX: %s -> %s", potx_path, converted_pptx_path)
        return self.template_converter.convert(
            potx_path=potx_path,
            pptx_path=converted_pptx_path,
            overwrite=True,
        )

    def _resolve_contract_path(self, contract_path: Optional[str | Path]) -> Path:
        return (
            Path(contract_path).expanduser().resolve()
            if contract_path is not None
            else self.engine_config.input.contract_path
        )

    def _get_slide_definitions(self, contract: Mapping[str, Any]) -> list[Mapping[str, Any]]:
        slides = contract.get("slides")

        if not isinstance(slides, list):
            raise DeckBuilderException("JSON contract must contain a 'slides' list.")

        for index, slide_definition in enumerate(slides):
            if not isinstance(slide_definition, Mapping):
                raise DeckBuilderException(f"Slide definition at index {index} must be an object.")

        return slides
