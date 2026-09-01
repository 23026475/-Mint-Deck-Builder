from __future__ import annotations

import logging
from dataclasses import dataclass
from typing import Any, Mapping, Sequence

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE

from presentation_engine.services.layout_policy import (
    MIN_CHART_HEIGHT as LAYOUT_MIN_CHART_HEIGHT,
    MIN_CHART_WIDTH as LAYOUT_MIN_CHART_WIDTH,
    normalize_chart_box,
)


logger = logging.getLogger(__name__)


class ChartHandlerException(Exception):
    """Raised when chart handling fails."""


class ChartPlaceholderNotFoundException(ChartHandlerException):
    """Raised when a chart placeholder cannot be found."""


class InvalidChartDataException(ChartHandlerException):
    """Raised when chart data cannot be normalised into chart categories and series."""


@dataclass(frozen=True)
class ChartRequest:
    archetype: str
    media_field: str
    source: str
    placeholder_idx: int
    chart_data: Any


@dataclass(frozen=True)
class NormalizedChartSeries:
    name: str
    values: list[float]


@dataclass(frozen=True)
class NormalizedChart:
    chart_type: str
    title: str | None
    categories: list[str]
    series: list[NormalizedChartSeries]

    @property
    def category_count(self) -> int:
        return len(self.categories)

    @property
    def series_count(self) -> int:
        return len(self.series)

    @property
    def is_empty(self) -> bool:
        return not self.categories or not self.series


@dataclass(frozen=True)
class PlaceholderGeometry:
    left: Any
    top: Any
    width: Any
    height: Any


class ChartHandler:
    """
    Handles chart insertion into existing PowerPoint chart placeholder geometry.

    The handler uses the configured placeholder as a coordinate source, removes
    the placeholder to prevent ghost prompt boxes, and adds a new chart shape at
    normalized chart-safe geometry.
    """

    # Public aliases used by unit tests and reports. Values are defined in layout_policy.
    MIN_CHART_WIDTH = LAYOUT_MIN_CHART_WIDTH
    MIN_CHART_HEIGHT = LAYOUT_MIN_CHART_HEIGHT

    DEFAULT_CHART_PLACEHOLDERS: dict[str, dict[str, int]] = {
        "chart": {
            "chart": 10,
        }
    }

    CHART_TYPE_MAP = {
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    }

    def __init__(
        self,
        placeholder_mappings: Mapping[str, Mapping[str, int]] | None = None,
        log: logging.Logger | None = None,
    ) -> None:
        self.placeholder_mappings = placeholder_mappings or self.DEFAULT_CHART_PLACEHOLDERS
        self.log = log or logger

    def process(self, slide: Any, slide_definition: Mapping[str, Any]) -> list[dict[str, Any]]:
        return self.insert_charts(slide, slide_definition)

    def insert_charts(self, slide: Any, slide_definition: Mapping[str, Any]) -> list[dict[str, Any]]:
        requests = self._collect_chart_requests(slide_definition)
        reports: list[dict[str, Any]] = []

        for request in requests:
            placeholder = self._placeholder_by_idx(slide).get(request.placeholder_idx)

            if placeholder is None:
                raise ChartPlaceholderNotFoundException(
                    f"Chart placeholder idx {request.placeholder_idx} was not found "
                    f"for archetype '{request.archetype}', media field '{request.media_field}'."
                )

            normalized = self._normalize_chart_data(request.chart_data)

            if normalized.is_empty:
                reports.append(
                    {
                        "archetype": request.archetype,
                        "media_field": request.media_field,
                        "source": request.source,
                        "placeholder_idx": request.placeholder_idx,
                        "status": "skipped_empty_chart",
                        "chart_type": normalized.chart_type,
                        "chart_title": normalized.title,
                        "category_count": normalized.category_count,
                        "series_count": normalized.series_count,
                    }
                )
                continue

            chart_shape = self._insert_chart(slide, placeholder, normalized)

            reports.append(
                {
                    "archetype": request.archetype,
                    "media_field": request.media_field,
                    "source": request.source,
                    "placeholder_idx": request.placeholder_idx,
                    "status": "inserted",
                    "chart_type": normalized.chart_type,
                    "chart_title": normalized.title,
                    "categories": normalized.categories,
                    "series": [
                        {
                            "name": series.name,
                            "values": series.values,
                        }
                        for series in normalized.series
                    ],
                    "category_count": normalized.category_count,
                    "series_count": normalized.series_count,
                    "shape_name": getattr(chart_shape, "name", None),
                    "geometry": {
                        "left": getattr(chart_shape, "left", None),
                        "top": getattr(chart_shape, "top", None),
                        "width": getattr(chart_shape, "width", None),
                        "height": getattr(chart_shape, "height", None),
                    },
                }
            )

        return reports

    def _collect_chart_requests(self, slide_definition: Mapping[str, Any]) -> list[ChartRequest]:
        archetype_value = slide_definition.get("archetype")

        if not isinstance(archetype_value, str):
            return []

        archetype = archetype_value.strip().lower()
        fields = slide_definition.get("fields", {})

        if not isinstance(fields, Mapping):
            return []

        chart_data = fields.get("chart")

        if chart_data is None:
            return []

        placeholder_idx = self._resolve_placeholder_idx(archetype, "chart")

        return [
            ChartRequest(
                archetype=archetype,
                media_field="chart",
                source="fields.chart",
                placeholder_idx=placeholder_idx,
                chart_data=chart_data,
            )
        ]

    def _resolve_placeholder_idx(self, archetype: str, media_field: str) -> int:
        archetype_mapping = self.placeholder_mappings.get(archetype)

        if not archetype_mapping:
            raise ChartHandlerException(
                f"No chart placeholder mapping found for archetype '{archetype}'."
            )

        if media_field not in archetype_mapping:
            raise ChartHandlerException(
                f"No chart placeholder mapping found for archetype '{archetype}', field '{media_field}'."
            )

        return int(archetype_mapping[media_field])

    def _normalize_chart_data(self, chart_data: Any) -> NormalizedChart:
        if chart_data is None:
            return NormalizedChart(chart_type="column", title=None, categories=[], series=[])

        if not isinstance(chart_data, Mapping):
            raise InvalidChartDataException("fields.chart must be an object.")

        chart_type = str(chart_data.get("type") or "column").strip().lower()

        if chart_type not in self.CHART_TYPE_MAP:
            raise InvalidChartDataException("fields.chart.type must be one of: column, bar.")

        title_value = chart_data.get("title")
        title = (
            str(title_value).strip()
            if title_value is not None and str(title_value).strip()
            else None
        )

        categories = chart_data.get("categories") or []
        series_payload = chart_data.get("series") or []

        if not isinstance(categories, Sequence) or isinstance(categories, (str, bytes)):
            raise InvalidChartDataException("fields.chart.categories must be a list of labels.")

        if not isinstance(series_payload, Sequence) or isinstance(series_payload, (str, bytes)):
            raise InvalidChartDataException("fields.chart.series must be a list of series objects.")

        normalized_categories = [str(value).strip() for value in categories]

        if not normalized_categories and not series_payload:
            return NormalizedChart(chart_type=chart_type, title=title, categories=[], series=[])

        if not normalized_categories:
            raise InvalidChartDataException("fields.chart.categories must contain at least one category.")

        normalized_series = self._normalize_series(series_payload, len(normalized_categories))

        if not normalized_series:
            return NormalizedChart(chart_type=chart_type, title=title, categories=[], series=[])

        return NormalizedChart(
            chart_type=chart_type,
            title=title,
            categories=normalized_categories,
            series=normalized_series,
        )

    def _normalize_series(
        self,
        series_payload: Sequence[Any],
        expected_value_count: int,
    ) -> list[NormalizedChartSeries]:
        normalized_series: list[NormalizedChartSeries] = []

        for index, item in enumerate(series_payload):
            if not isinstance(item, Mapping):
                raise InvalidChartDataException("Each chart series must be an object.")

            name = str(item.get("name") or f"Series {index + 1}").strip()
            values = item.get("values")

            if not isinstance(values, Sequence) or isinstance(values, (str, bytes)):
                raise InvalidChartDataException("Each chart series must contain a values list.")

            numeric_values = [self._coerce_numeric_value(value) for value in values]

            if len(numeric_values) != expected_value_count:
                raise InvalidChartDataException(
                    "Each chart series must contain the same number of values as fields.chart.categories."
                )

            normalized_series.append(NormalizedChartSeries(name=name, values=numeric_values))

        return normalized_series

    def _coerce_numeric_value(self, value: Any) -> float:
        if isinstance(value, bool):
            raise InvalidChartDataException("Chart values must be numeric and cannot be booleans.")

        if isinstance(value, (int, float)):
            return float(value)

        if isinstance(value, str) and value.strip():
            try:
                return float(value.strip())
            except ValueError as exc:
                raise InvalidChartDataException(
                    f"Chart value is not numeric: {value!r}."
                ) from exc

        raise InvalidChartDataException(f"Chart value is not numeric: {value!r}.")

    def _insert_chart(self, slide: Any, placeholder: Any, chart: NormalizedChart) -> Any:
        """Insert a visible chart shape using the placeholder as coordinate source."""
        shapes = getattr(slide, "shapes", None)

        if shapes is None or not hasattr(shapes, "add_chart"):
            raise ChartPlaceholderNotFoundException("Slide does not support shapes.add_chart().")

        geometry = self._resolve_placeholder_geometry(slide, placeholder)
        slide_width, slide_height = self._slide_dimensions(slide)
        box = normalize_chart_box(
            geometry.left,
            geometry.top,
            geometry.width,
            geometry.height,
            slide_width=slide_width,
            slide_height=slide_height,
        )

        # Important: remove the prompt placeholder before adding the replacement
        # chart frame. This avoids a ghost placeholder or prompt remaining under
        # the chart after generation.
        self._remove_placeholder(placeholder)

        chart_data = self._to_category_chart_data(chart)
        chart_shape = shapes.add_chart(
            self.CHART_TYPE_MAP[chart.chart_type],
            box.left,
            box.top,
            box.width,
            box.height,
            chart_data,
        )

        if chart.title and hasattr(chart_shape, "chart"):
            try:
                chart_shape.chart.has_title = True
                chart_shape.chart.chart_title.text_frame.text = chart.title
            except (AttributeError, TypeError, ValueError):
                self.log.debug("Could not set chart title", exc_info=True)

        self._apply_minimal_visible_chart_formatting(chart_shape)
        return chart_shape

    def _resolve_placeholder_geometry(self, slide: Any, placeholder: Any) -> PlaceholderGeometry:
        left = getattr(placeholder, "left", None)
        top = getattr(placeholder, "top", None)
        width = getattr(placeholder, "width", None)
        height = getattr(placeholder, "height", None)

        if left is not None and top is not None and width is not None and height is not None:
            return PlaceholderGeometry(left=left, top=top, width=width, height=height)

        layout_placeholder = self._matching_layout_placeholder(slide, placeholder)
        if layout_placeholder is not None:
            left = left if left is not None else getattr(layout_placeholder, "left", None)
            top = top if top is not None else getattr(layout_placeholder, "top", None)
            width = width if width is not None else getattr(layout_placeholder, "width", None)
            height = height if height is not None else getattr(layout_placeholder, "height", None)

        return PlaceholderGeometry(left=left, top=top, width=width, height=height)

    def _matching_layout_placeholder(self, slide: Any, placeholder: Any) -> Any | None:
        idx = self._placeholder_idx(placeholder)
        slide_layout = getattr(slide, "slide_layout", None)
        layout_placeholders = getattr(slide_layout, "placeholders", []) if slide_layout is not None else []

        for candidate in layout_placeholders:
            if self._placeholder_idx(candidate) == idx:
                return candidate

        return None

    def _slide_dimensions(self, slide: Any) -> tuple[Any | None, Any | None]:
        """Best-effort slide dimension lookup without changing handler API."""
        try:
            presentation = slide.part.package.presentation_part.presentation
            return presentation.slide_width, presentation.slide_height
        except AttributeError:
            return None, None

    def _apply_minimal_visible_chart_formatting(self, chart_shape: Any) -> None:
        """Apply minimal explicit fill/outline styling for visible bars/columns."""
        chart_object = getattr(chart_shape, "chart", None)

        if chart_object is None:
            return

        try:
            chart_object.has_legend = False
        except (AttributeError, TypeError, ValueError):
            self.log.debug("Could not disable chart legend", exc_info=True)

        series_colours = [
            RGBColor(27, 95, 174),
            RGBColor(0, 166, 166),
            RGBColor(126, 232, 201),
        ]

        try:
            for index, series in enumerate(chart_object.series):
                colour = series_colours[index % len(series_colours)]
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = colour
                series.format.line.color.rgb = colour
        except (AttributeError, TypeError, ValueError):
            self.log.debug("Could not apply visible chart series formatting", exc_info=True)

    def _remove_placeholder(self, placeholder: Any) -> None:
        element = getattr(placeholder, "_element", None)

        if element is None:
            return

        try:
            parent = element.getparent()
        except AttributeError:
            return

        if parent is not None:
            try:
                parent.remove(element)
            except (AttributeError, TypeError, ValueError):
                self.log.debug("Could not remove chart placeholder", exc_info=True)

    def _to_category_chart_data(self, chart: NormalizedChart) -> CategoryChartData:
        chart_data = CategoryChartData()
        chart_data.categories = chart.categories

        for series in chart.series:
            chart_data.add_series(series.name, series.values)

        return chart_data

    def _placeholder_by_idx(self, slide: Any) -> dict[int, Any]:
        placeholders: dict[int, Any] = {}

        for placeholder in getattr(slide, "placeholders", []):
            idx = self._placeholder_idx(placeholder)

            if idx is not None:
                placeholders[idx] = placeholder

        return placeholders

    def _placeholder_idx(self, shape: Any) -> int | None:
        try:
            return int(shape.placeholder_format.idx)
        except (AttributeError, TypeError, ValueError):
            return None
