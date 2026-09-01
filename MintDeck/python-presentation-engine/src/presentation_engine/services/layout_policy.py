"""Small layout-policy helpers for the Python Presentation Engine.

Purpose:
- Keep layout behaviour consistent across existing handlers.
- Enforce minimum component sizes for media objects such as charts.
- Clamp normalized component boxes inside known slide bounds.
- Apply predictable text-frame behaviour without replacing typography from
  the approved PowerPoint template.

This module is intentionally small. It is not a replacement layout engine.
"""

from __future__ import annotations

import logging
from dataclasses import dataclass
from typing import Any

from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt


logger = logging.getLogger(__name__)


@dataclass(frozen=True)
class Box:
    """A positioned rectangular region expressed in PowerPoint EMUs."""

    left: Any
    top: Any
    width: Any
    height: Any


MIN_CHART_WIDTH = Inches(5.5)
MIN_CHART_HEIGHT = Inches(3.0)
MIN_TEXT_FONT = Pt(14)
DEFAULT_TEXT_FONT = Pt(18)


def normalize_chart_box(
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    *,
    min_width: Any = MIN_CHART_WIDTH,
    min_height: Any = MIN_CHART_HEIGHT,
    slide_width: Any | None = None,
    slide_height: Any | None = None,
) -> Box:
    """Return a chart-safe box using a placeholder as coordinate source.

    Width and height are expanded to the configured minimums. When slide
    dimensions are supplied, the resulting box is clamped inside the slide
    canvas instead of growing beyond the right or bottom edge.
    """

    normalized_width = _at_least(width, min_width)
    normalized_height = _at_least(height, min_height)

    normalized_left = _safe_coordinate(left, 0)
    normalized_top = _safe_coordinate(top, 0)

    if slide_width is not None:
        normalized_width = _at_most(normalized_width, slide_width)
        normalized_left = _clamp_coordinate(
            normalized_left,
            minimum=0,
            maximum=max(0, slide_width - normalized_width),
        )
    else:
        normalized_left = _at_least(normalized_left, 0)

    if slide_height is not None:
        normalized_height = _at_most(normalized_height, slide_height)
        normalized_top = _clamp_coordinate(
            normalized_top,
            minimum=0,
            maximum=max(0, slide_height - normalized_height),
        )
    else:
        normalized_top = _at_least(normalized_top, 0)

    return Box(
        left=normalized_left,
        top=normalized_top,
        width=normalized_width,
        height=normalized_height,
    )


def configure_text_frame(
    shape: Any,
    *,
    font_size: Any | None = None,
    min_font_size: Any | None = None,
) -> None:
    """Configure fixed-box text behaviour while preserving template typography.

    Default behaviour:
    - enable word wrapping;
    - use top vertical alignment;
    - apply small, consistent text-frame margins;
    - turn off PowerPoint auto-fit intent;
    - preserve font family, size, colour, weight and style inherited from the
      approved template.

    Font sizes are changed only if a caller explicitly supplies ``font_size``
    or ``min_font_size``. Calling ``configure_text_frame(shape)`` therefore
    does not flatten the template's typography hierarchy.

    Lightweight fake shapes used by unit tests are safely ignored when they do
    not expose a real ``text_frame`` object.
    """

    if not getattr(shape, "has_text_frame", False):
        return

    text_frame = getattr(shape, "text_frame", None)
    if text_frame is None:
        return

    shape_name = getattr(shape, "name", "unknown")

    try:
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        text_frame.margin_left = Inches(0.08)
        text_frame.margin_right = Inches(0.08)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)

        # Do not rely on PowerPoint to calculate a shrink factor after the deck
        # has been generated. Content ceilings are validated before rendering.
        text_frame.auto_size = None
    except (AttributeError, TypeError, ValueError):
        logger.debug(
            "Could not configure text frame for shape %s",
            shape_name,
            exc_info=True,
        )
        return

    # The approved template remains the source of truth for typography. Apply
    # size constraints only when a specific caller explicitly requests them.
    if font_size is None and min_font_size is None:
        return

    paragraphs = list(getattr(text_frame, "paragraphs", []) or [])

    for paragraph in paragraphs:
        paragraph_font = getattr(paragraph, "font", None)
        if paragraph_font is not None:
            target_size = _resolve_optional_font_size(
                current_size=getattr(paragraph_font, "size", None),
                max_size=font_size,
                minimum_size=min_font_size,
            )

            if target_size is not None:
                try:
                    paragraph_font.size = target_size
                except (AttributeError, TypeError, ValueError):
                    logger.debug(
                        "Could not set paragraph font size for shape %s",
                        shape_name,
                        exc_info=True,
                    )

        for run in list(getattr(paragraph, "runs", []) or []):
            run_font = getattr(run, "font", None)
            if run_font is None:
                continue

            target_size = _resolve_optional_font_size(
                current_size=getattr(run_font, "size", None),
                max_size=font_size,
                minimum_size=min_font_size,
            )

            if target_size is None:
                continue

            try:
                run_font.size = target_size
            except (AttributeError, TypeError, ValueError):
                logger.debug(
                    "Could not set run font size for shape %s",
                    shape_name,
                    exc_info=True,
                )


def _resolve_optional_font_size(
    *,
    current_size: Any | None,
    max_size: Any | None,
    minimum_size: Any | None,
) -> Any | None:
    """Resolve an explicitly requested font-size constraint.

    ``None`` means the existing or inherited template font size must remain
    unchanged.
    """

    if current_size is None:
        if max_size is None:
            return None
        current_size = max_size

    target_size = current_size

    try:
        if max_size is not None and target_size > max_size:
            target_size = max_size

        if minimum_size is not None and target_size < minimum_size:
            target_size = minimum_size
    except (TypeError, ValueError):
        logger.debug("Could not resolve optional font size", exc_info=True)
        return max_size

    return target_size


def _safe_coordinate(value: Any, fallback: Any) -> Any:
    """Return a usable coordinate, falling back when the value is missing."""

    return fallback if value is None else value


def _at_least(value: Any, minimum: Any) -> Any:
    """Return ``value`` constrained to a minimum."""

    if value is None:
        return minimum

    try:
        return value if value >= minimum else minimum
    except (TypeError, ValueError):
        logger.debug("Could not compare layout value against minimum", exc_info=True)
        return minimum


def _at_most(value: Any, maximum: Any) -> Any:
    """Return ``value`` constrained to a maximum."""

    if value is None:
        return maximum

    try:
        return value if value <= maximum else maximum
    except (TypeError, ValueError):
        logger.debug("Could not compare layout value against maximum", exc_info=True)
        return maximum


def _clamp_coordinate(value: Any, *, minimum: Any, maximum: Any) -> Any:
    """Clamp a coordinate between inclusive minimum and maximum values."""

    try:
        return max(minimum, min(value, maximum))
    except (TypeError, ValueError):
        logger.debug("Could not clamp layout coordinate", exc_info=True)
        return minimum


def _clamp_font_size(value: Any, minimum: Any, maximum: Any) -> Any:
    """Backward-compatible explicit font-size clamp helper."""

    size = value if value is not None else maximum

    try:
        return max(minimum, min(size, maximum))
    except (TypeError, ValueError):
        logger.debug("Could not clamp font size", exc_info=True)
        return maximum
