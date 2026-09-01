"""Deterministic text-container sizing for approved PowerPoint placeholders.

The contract remains the source of visible wording and the PowerPoint template
remains the source of typography, colour, theme, and standard geometry. This
policy adjusts only the height of placeholders explicitly configured as
``dynamic``.
"""

from __future__ import annotations

import json
import logging
import math
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Mapping

from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches


logger = logging.getLogger(__name__)
EMU_PER_INCH = 914400


class TextBoxLayoutError(Exception):
    """Raised when text-box layout configuration is missing or invalid."""


@dataclass(frozen=True)
class TextBoxRule:
    archetype: str
    placeholder_name: str
    placeholder_idx: int
    mode: str
    min_height: float | None = None
    max_height: float | None = None


@dataclass(frozen=True)
class TextBoxMeasurement:
    estimated_lines: int
    estimated_height_inches: float
    applied_height_inches: float


class TextBoxLayoutPolicy:
    """Load and apply text-container sizing rules by archetype and field name."""

    DEFAULT_CONFIG_PATH = Path("config/text-box-layout-policy.json")

    def __init__(
        self,
        config_path: str | Path | None = None,
        config_data: Mapping[str, Any] | None = None,
    ) -> None:
        self.config_path = Path(config_path) if config_path else self.DEFAULT_CONFIG_PATH
        self._config = dict(config_data) if config_data is not None else self._load()
        self._validate()

    def _load(self) -> dict[str, Any]:
        if not self.config_path.exists():
            raise TextBoxLayoutError(
                f"Text-box layout configuration does not exist: {self.config_path}"
            )

        try:
            value = json.loads(self.config_path.read_text(encoding="utf8"))
        except (OSError, json.JSONDecodeError) as exc:
            raise TextBoxLayoutError(
                f"Could not load text-box layout configuration: {self.config_path}"
            ) from exc

        if not isinstance(value, dict):
            raise TextBoxLayoutError(
                "Text-box layout configuration must contain a JSON object."
            )

        return value

    def _validate(self) -> None:
        if not isinstance(self._config.get("defaults"), Mapping):
            raise TextBoxLayoutError("Configuration must contain a defaults object.")

        if not isinstance(self._config.get("archetypes"), Mapping):
            raise TextBoxLayoutError("Configuration must contain an archetypes object.")

        for archetype, rules in self._config["archetypes"].items():
            if not isinstance(rules, Mapping):
                raise TextBoxLayoutError(
                    f"Rules for archetype '{archetype}' must be an object."
                )

            for name, rule in rules.items():
                if not isinstance(rule, Mapping):
                    raise TextBoxLayoutError(
                        f"Rule '{archetype}.{name}' must be an object."
                    )

                if not isinstance(rule.get("idx"), int):
                    raise TextBoxLayoutError(
                        f"Rule '{archetype}.{name}.idx' must be an integer."
                    )

                mode = rule.get("mode", self._config["defaults"].get("mode"))
                if mode not in {"fixed", "dynamic"}:
                    raise TextBoxLayoutError(
                        f"Rule '{archetype}.{name}.mode' must be fixed or dynamic."
                    )

    def resolve(self, archetype: str, placeholder_name: str) -> TextBoxRule | None:
        archetype_key = (archetype or "").strip().lower()
        rules = self._config["archetypes"].get(archetype_key)
        if not isinstance(rules, Mapping):
            return None

        payload = rules.get(placeholder_name)
        if not isinstance(payload, Mapping):
            return None

        return TextBoxRule(
            archetype=archetype_key,
            placeholder_name=placeholder_name,
            placeholder_idx=int(payload["idx"]),
            mode=str(payload.get("mode", self._config["defaults"].get("mode", "fixed"))),
            min_height=_optional_float(payload.get("min_height")),
            max_height=_optional_float(payload.get("max_height")),
        )

    def apply(
        self,
        shape: Any,
        *,
        archetype: str,
        placeholder_name: str,
        text: str,
        font_size_pt: float | None = None,
    ) -> TextBoxMeasurement | None:
        rule = self.resolve(archetype, placeholder_name)
        if rule is None or rule.mode != "dynamic":
            return None

        if not getattr(shape, "has_text_frame", False):
            return None

        text_frame = getattr(shape, "text_frame", None)
        if text_frame is None:
            return None

        # Resolve the real placeholder's complete effective geometry before any
        # mutation. For inherited placeholders, these values are supplied by the
        # matching layout placeholder through python-pptx's property resolution.
        effective_left = getattr(shape, "left", None)
        effective_top = getattr(shape, "top", None)
        effective_width = getattr(shape, "width", None)
        effective_height = getattr(shape, "height", None)

        if None in (
            effective_left,
            effective_top,
            effective_width,
            effective_height,
        ):
            logger.debug(
                "Could not resolve complete geometry for dynamic placeholder %s.%s",
                archetype,
                placeholder_name,
            )
            return None

        width_inches = _emu_to_inches(effective_width)
        if width_inches <= 0:
            return None

        resolved_font_size = font_size_pt or _font_size_from_shape(shape) or 14.0
        defaults = self._config["defaults"]
        line_height_multiplier = float(defaults.get("line_height_multiplier", 1.2))
        paragraph_gap_lines = float(defaults.get("paragraph_gap_lines", 0.35))
        chars_constant = float(defaults.get("characters_per_point_inch", 27.8))

        margin_left = float(defaults.get("margin_left", 0.08))
        margin_right = float(defaults.get("margin_right", 0.08))
        margin_top = float(defaults.get("margin_top", 0.05))
        margin_bottom = float(defaults.get("margin_bottom", 0.05))
        usable_width = max(0.1, width_inches - margin_left - margin_right)

        chars_per_inch = max(8.0, chars_constant * 14.0 / resolved_font_size)
        chars_per_line = max(1, int(usable_width * chars_per_inch))
        paragraphs = (text or "").splitlines() or [""]

        estimated_lines = 0
        for paragraph in paragraphs:
            normalized = " ".join(paragraph.split())
            estimated_lines += max(1, math.ceil(len(normalized) / chars_per_line))

        if len(paragraphs) > 1:
            estimated_lines += math.ceil(
                (len(paragraphs) - 1) * paragraph_gap_lines
            )

        line_height_inches = (resolved_font_size / 72.0) * line_height_multiplier
        estimated_height = (
            estimated_lines * line_height_inches
            + margin_top
            + margin_bottom
        )

        applied_height = estimated_height
        if rule.min_height is not None:
            applied_height = max(applied_height, rule.min_height)
        if rule.max_height is not None:
            applied_height = min(applied_height, rule.max_height)

        try:
            # Materialize the complete effective transform. Writing height alone
            # on an inherited placeholder creates an ext element with cx=0.
            shape.left = effective_left
            shape.top = effective_top
            shape.width = effective_width
            shape.height = Inches(applied_height)
            text_frame.vertical_anchor = MSO_ANCHOR.TOP
        except (AttributeError, TypeError, ValueError):
            logger.debug(
                "Could not resize text box %s for %s.%s",
                getattr(shape, "name", "unknown"),
                archetype,
                placeholder_name,
                exc_info=True,
            )
            return None

        return TextBoxMeasurement(
            estimated_lines=estimated_lines,
            estimated_height_inches=estimated_height,
            applied_height_inches=applied_height,
        )


def _font_size_from_shape(shape: Any) -> float | None:
    text_frame = getattr(shape, "text_frame", None)
    if text_frame is None:
        return None

    for paragraph in list(getattr(text_frame, "paragraphs", []) or []):
        paragraph_font = getattr(paragraph, "font", None)
        paragraph_size = (
            getattr(paragraph_font, "size", None)
            if paragraph_font is not None
            else None
        )
        if paragraph_size is not None:
            return float(paragraph_size.pt)

        for run in list(getattr(paragraph, "runs", []) or []):
            run_font = getattr(run, "font", None)
            run_size = (
                getattr(run_font, "size", None)
                if run_font is not None
                else None
            )
            if run_size is not None:
                return float(run_size.pt)

    return None


def _emu_to_inches(value: Any) -> float:
    if value is None:
        return 0.0

    try:
        return float(value) / EMU_PER_INCH
    except (TypeError, ValueError):
        return 0.0


def _optional_float(value: Any) -> float | None:
    if value is None:
        return None

    try:
        return float(value)
    except (TypeError, ValueError) as exc:
        raise TextBoxLayoutError(
            f"Expected a numeric layout value, received {value!r}."
        ) from exc
